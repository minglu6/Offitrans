#!/usr/bin/env python3
"""
Batch Processing Examples for Offitrans

This example demonstrates how to process multiple files efficiently
with batch operations, parallel processing, and progress monitoring.
"""

import os
import time
import threading
from pathlib import Path
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import List, Dict, Any

from offitrans import ExcelProcessor, GoogleTranslator
from offitrans.processors import get_processor_by_extension
from offitrans.core.config import Config

# Set up logging
import logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)


class BatchProcessor:
    """
    Enhanced batch processor for handling multiple files
    """
    
    def __init__(self, config: Config = None):
        """
        Initialize batch processor
        
        Args:
            config: Configuration for processors and translators
        """
        self.config = config or Config()
        self.results = []
        self.lock = threading.Lock()
        
        # Statistics
        self.stats = {
            'total_files': 0,
            'successful_files': 0,
            'failed_files': 0,
            'start_time': None,
            'end_time': None,
            'total_processing_time': 0
        }
    
    def process_file(self, input_file: str, output_file: str, target_language: str = "en") -> Dict[str, Any]:
        """
        Process a single file with error handling and timing
        
        Args:
            input_file: Path to input file
            output_file: Path to output file
            target_language: Target language code
            
        Returns:
            Processing result dictionary
        """
        start_time = time.time()
        result = {
            'input_file': input_file,
            'output_file': output_file,
            'target_language': target_language,
            'success': False,
            'error': None,
            'processing_time': 0,
            'file_size': 0,
            'texts_translated': 0
        }
        
        try:
            # Get file size
            if os.path.exists(input_file):
                result['file_size'] = os.path.getsize(input_file)
            
            # Get appropriate processor
            processor = get_processor_by_extension(input_file, config=self.config)
            
            # Process the file
            success = processor.process_file(input_file, output_file, target_language)
            
            if success:
                result['success'] = True
                
                # Get processor statistics
                proc_stats = processor.get_stats()
                result['texts_translated'] = proc_stats.get('total_texts_translated', 0)
                
                logger.info(f"Successfully processed: {input_file}")
            else:
                result['error'] = "Processing failed"
                logger.error(f"Failed to process: {input_file}")
                
        except Exception as e:
            result['error'] = str(e)
            logger.error(f"Error processing {input_file}: {e}")
        
        # Calculate processing time
        result['processing_time'] = time.time() - start_time
        
        # Update statistics thread-safely
        with self.lock:
            self.stats['total_files'] += 1
            if result['success']:
                self.stats['successful_files'] += 1
            else:
                self.stats['failed_files'] += 1
            
            self.results.append(result)
        
        return result
    
    def process_files_sequential(self, file_pairs: List[tuple], target_language: str = "en") -> List[Dict[str, Any]]:
        """
        Process files sequentially (one after another)
        
        Args:
            file_pairs: List of (input_file, output_file) tuples
            target_language: Target language code
            
        Returns:
            List of processing results
        """
        print(f"Processing {len(file_pairs)} files sequentially...")
        
        self.stats['start_time'] = time.time()
        
        for i, (input_file, output_file) in enumerate(file_pairs, 1):
            print(f"\nProcessing file {i}/{len(file_pairs)}: {input_file}")
            self.process_file(input_file, output_file, target_language)
        
        self.stats['end_time'] = time.time()
        self.stats['total_processing_time'] = self.stats['end_time'] - self.stats['start_time']
        
        return self.results
    
    def process_files_parallel(self, file_pairs: List[tuple], target_language: str = "en", max_workers: int = 3) -> List[Dict[str, Any]]:
        """
        Process files in parallel using multiple threads
        
        Args:
            file_pairs: List of (input_file, output_file) tuples
            target_language: Target language code
            max_workers: Maximum number of parallel workers
            
        Returns:
            List of processing results
        """
        print(f"Processing {len(file_pairs)} files in parallel with {max_workers} workers...")
        
        self.stats['start_time'] = time.time()
        
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            # Submit all tasks
            future_to_file = {
                executor.submit(self.process_file, input_file, output_file, target_language): (input_file, output_file)
                for input_file, output_file in file_pairs
            }
            
            # Process completed tasks
            for future in as_completed(future_to_file):
                input_file, output_file = future_to_file[future]
                try:
                    result = future.result()
                    status = "OK" if result['success'] else "FAILED"
                    print(f"{status} {input_file} -> {output_file} ({result['processing_time']:.1f}s)")
                except Exception as e:
                    print(f"FAILED {input_file} -> Exception: {e}")
        
        self.stats['end_time'] = time.time()
        self.stats['total_processing_time'] = self.stats['end_time'] - self.stats['start_time']
        
        return self.results
    
    def print_summary(self):
        """Print processing summary and statistics"""
        print("\n" + "=" * 60)
        print("Batch Processing Summary")
        print("=" * 60)
        
        print(f"Total files processed: {self.stats['total_files']}")
        print(f"Successful: {self.stats['successful_files']}")
        print(f"Failed: {self.stats['failed_files']}")
        print(f"Success rate: {(self.stats['successful_files'] / max(1, self.stats['total_files']) * 100):.1f}%")
        print(f"Total processing time: {self.stats['total_processing_time']:.1f} seconds")
        
        if self.stats['successful_files'] > 0:
            avg_time = self.stats['total_processing_time'] / self.stats['successful_files']
            print(f"Average time per file: {avg_time:.1f} seconds")
        
        # Show failed files
        failed_files = [r for r in self.results if not r['success']]
        if failed_files:
            print(f"\nFailed files:")
            for result in failed_files:
                print(f"   - {result['input_file']}: {result['error']}")
        
        # Show successful files
        successful_files = [r for r in self.results if r['success']]
        if successful_files:
            print(f"\nSuccessfully processed files:")
            for result in successful_files:
                print(f"   - {result['input_file']} ({result['processing_time']:.1f}s, {result['texts_translated']} texts)")


def create_sample_files():
    """Create sample files for batch processing demo"""
    sample_dir = Path("examples/sample_files/batch_demo")
    sample_dir.mkdir(parents=True, exist_ok=True)
    
    # Create sample Excel files
    try:
        from openpyxl import Workbook
        
        sample_data = [
            ("sample1.xlsx", ["你好世界", "这是第一个测试文件", "包含中文内容"]),
            ("sample2.xlsx", ["欢迎使用", "批量处理功能", "非常方便"]),
            ("sample3.xlsx", ["测试文件", "包含多种内容", "数字123", "邮箱test@example.com"]),
            ("sample4.xlsx", ["最后一个", "测试文件", "批量处理完成"])
        ]
        
        for filename, texts in sample_data:
            file_path = sample_dir / filename
            
            wb = Workbook()
            ws = wb.active
            ws.title = "测试数据"
            
            for i, text in enumerate(texts, 1):
                ws[f'A{i}'] = text
            
            wb.save(file_path)
            print(f"Created sample file: {file_path}")
        
    except ImportError:
        print("openpyxl not available, cannot create Excel sample files")
        return []
    
    return [sample_dir / f for f in ["sample1.xlsx", "sample2.xlsx", "sample3.xlsx", "sample4.xlsx"]]


def demo_sequential_processing():
    """Demo sequential file processing"""
    print("=" * 60)
    print("Sequential Processing Demo")
    print("=" * 60)
    
    # Create sample files
    sample_files = create_sample_files()
    if not sample_files:
        return
    
    # Prepare file pairs
    file_pairs = []
    for input_file in sample_files:
        output_file = input_file.parent / f"{input_file.stem}_translated_seq{input_file.suffix}"
        file_pairs.append((str(input_file), str(output_file)))
    
    # Create processor with conservative settings for sequential processing
    config = Config()
    config.translator.max_workers = 2
    config.cache.enabled = True
    
    processor = BatchProcessor(config)
    
    # Process files sequentially
    results = processor.process_files_sequential(file_pairs, target_language="en")
    
    # Print summary
    processor.print_summary()


def demo_parallel_processing():
    """Demo parallel file processing"""
    print("\n" + "=" * 60)
    print("Parallel Processing Demo")
    print("=" * 60)
    
    # Create sample files
    sample_files = create_sample_files()
    if not sample_files:
        return
    
    # Prepare file pairs
    file_pairs = []
    for input_file in sample_files:
        output_file = input_file.parent / f"{input_file.stem}_translated_par{input_file.suffix}"
        file_pairs.append((str(input_file), str(output_file)))
    
    # Create processor with settings optimized for parallel processing
    config = Config()
    config.translator.max_workers = 1  # Reduce per-processor workers for parallel processing
    config.cache.enabled = True
    
    processor = BatchProcessor(config)
    
    # Process files in parallel
    results = processor.process_files_parallel(file_pairs, target_language="en", max_workers=3)
    
    # Print summary
    processor.print_summary()


def demo_mixed_file_types():
    """Demo processing different file types in a batch"""
    print("\n" + "=" * 60)
    print("Mixed File Types Processing Demo")
    print("=" * 60)
    
    sample_dir = Path("examples/sample_files/mixed_demo")
    sample_dir.mkdir(parents=True, exist_ok=True)
    
    # Create different types of sample files
    sample_files = []
    
    # Excel file
    try:
        from openpyxl import Workbook
        excel_file = sample_dir / "mixed_excel.xlsx"
        wb = Workbook()
        ws = wb.active
        ws['A1'] = "Excel文件测试"
        ws['A2'] = "包含中文内容"
        wb.save(excel_file)
        sample_files.append(excel_file)
        print(f"Created Excel sample: {excel_file}")
    except ImportError:
        pass
    
    # Word file
    try:
        from docx import Document
        word_file = sample_dir / "mixed_word.docx"
        doc = Document()
        doc.add_paragraph("Word文档测试")
        doc.add_paragraph("这是中文段落")
        doc.save(word_file)
        sample_files.append(word_file)
        print(f"Created Word sample: {word_file}")
    except ImportError:
        pass
    
    # PowerPoint file
    try:
        from pptx import Presentation
        ppt_file = sample_dir / "mixed_ppt.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "PPT演示文稿"
        slide.shapes.placeholders[1].text = "包含中文内容的幻灯片"
        prs.save(ppt_file)
        sample_files.append(ppt_file)
        print(f"Created PowerPoint sample: {ppt_file}")
    except ImportError:
        pass
    
    if not sample_files:
        print("No sample files could be created (missing dependencies)")
        return
    
    # Prepare file pairs
    file_pairs = []
    for input_file in sample_files:
        output_file = input_file.parent / f"{input_file.stem}_translated{input_file.suffix}"
        file_pairs.append((str(input_file), str(output_file)))
    
    # Process mixed file types
    config = Config()
    processor = BatchProcessor(config)
    
    results = processor.process_files_parallel(file_pairs, target_language="en", max_workers=2)
    
    # Print summary
    processor.print_summary()


def demo_progress_monitoring():
    """Demo with progress monitoring and real-time updates"""
    print("\n" + "=" * 60)
    print("Progress Monitoring Demo")
    print("=" * 60)
    
    # This would be a more sophisticated version with progress bars
    # For now, we'll show the concept with simple progress updates
    
    sample_files = create_sample_files()
    if not sample_files:
        return
    
    print("Processing with progress monitoring...")
    
    config = Config()
    processor = BatchProcessor(config)
    
    # Simulate progress monitoring
    for i, input_file in enumerate(sample_files, 1):
        output_file = input_file.parent / f"{input_file.stem}_monitored{input_file.suffix}"
        
        print(f"\nProgress: {i}/{len(sample_files)} ({(i/len(sample_files)*100):.1f}%)")
        print(f"Current file: {input_file.name}")
        
        start_time = time.time()
        result = processor.process_file(str(input_file), str(output_file), "en")
        end_time = time.time()
        
        print(f"   Status: {'Success' if result['success'] else 'Failed'}")
        print(f"   Time: {end_time - start_time:.1f}s")
        if result['texts_translated']:
            print(f"   Texts translated: {result['texts_translated']}")
    
    processor.print_summary()


def main():
    """
    Main function to run all batch processing examples
    """
    print("Offitrans Batch Processing Examples")
    print("This example demonstrates efficient batch processing capabilities")
    
    # Run all demos
    demo_sequential_processing()
    demo_parallel_processing()
    demo_mixed_file_types()
    demo_progress_monitoring()
    
    print("\n" + "=" * 60)
    print("Batch processing examples completed!")
    print("=" * 60)
    print("Performance Tips:")
    print("   1. Use parallel processing for multiple independent files")
    print("   2. Reduce max_workers per processor when using parallel file processing")
    print("   3. Enable caching to avoid re-translating the same content")
    print("   4. Monitor memory usage with large files or many parallel workers")
    print("   5. Consider file size and complexity when setting max_workers")
    print("\nCheck the generated files in examples/sample_files/")


if __name__ == "__main__":
    main()