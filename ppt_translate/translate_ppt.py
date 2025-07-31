import logging
import os
import re
import sys
import xml.etree.ElementTree as ET
import zipfile
from typing import List, Dict

# 设置日志级别，只显示INFO及以上级别的消息
logging.basicConfig(level=logging.INFO, format='%(levelname)s:%(name)s:%(message)s')

# 添加父目录到路径以导入翻译工具
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt
from translate_tools.base import Translator
from translate_tools.sungrow_translate import SunTranslator
from translate_tools.utils import should_translate, normalize_text


def _clean_smartart_text(text: str) -> str:
    """清理和格式化SmartArt文本"""
    if not text:
        return ""

    # 移除多余的空格
    text = ' '.join(text.split())

    # 过滤掉太短或无意义的文本
    if (len(text) < 3 or
            text.isdigit() or
            text in ['true', 'false', 'exact', 'lvl']):
        return ""

    # 过滤掉纯标点或纯数字加标点的文本
    import re
    if re.match(r'^[\d\s.,:\-()＜＞]+$', text):
        return ""

    return text


def _extract_smartart_xml_text(root, slide_idx: int, shape_idx: int, shape, text_data: List[Dict]):
    """从SmartArt XML中提取文本"""
    try:
        text_elements = set()

        # 查找所有可能包含文本的元素
        for elem in root.iter():
            # 检查元素文本
            if elem.text and elem.text.strip():
                text = elem.text.strip()
                if (len(text) > 1 and
                        not text.startswith('{') and
                        not text.isdigit() and
                        text not in text_elements):
                    text_elements.add(text)
                    text_data.append({
                        'slide_index': slide_idx,
                        'shape_index': shape_idx,
                        'text': text,
                        'smartart_obj': shape,
                        'xml_element': elem,
                        'is_smartart': True,
                        'is_xml_text': True
                    })

            # 检查元素属性中的文本
            for attr_name, attr_value in elem.attrib.items():
                if (attr_name in ['val', 'text', 'title', 'name'] and
                        attr_value and attr_value.strip() and
                        len(attr_value.strip()) > 2 and
                        not attr_value.isdigit() and
                        attr_value.strip() not in text_elements):
                    text_elements.add(attr_value.strip())
                    text_data.append({
                        'slide_index': slide_idx,
                        'shape_index': shape_idx,
                        'text': attr_value.strip(),
                        'smartart_obj': shape,
                        'xml_element': elem,
                        'is_smartart': True,
                        'is_xml_attr': True
                    })

    except Exception as e:
        logging.warning(f"[extract_smartart_xml_text] SmartArt XML文本提取失败: {e}")


def _extract_smartart_runs(shape, slide_idx: int, shape_idx: int, text_data: List[Dict]):
    """从SmartArt的文本运行中提取文本"""
    try:
        if hasattr(shape, 'text_frame') and shape.text_frame:
            for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                for run_idx, run in enumerate(paragraph.runs):
                    if run.text.strip():
                        text_data.append({
                            'slide_index': slide_idx,
                            'shape_index': shape_idx,
                            'paragraph_index': para_idx,
                            'run_index': run_idx,
                            'text': run.text.strip(),
                            'run_obj': run,
                            'is_smartart': True,
                            'is_run': True
                        })
    except Exception as e:
        logging.warning(f"[extract_smartart_runs] SmartArt runs提取失败: {e}")


def _extract_from_media_shape(text_data: List[Dict], slide_idx: int, shape_idx: int, shape):
    """从图片和媒体形状中提取文本"""
    try:
        # 图片可能有关联的文本框或标题
        if hasattr(shape, 'element'):
            # 查找图片的替代文本或标题
            xml_text = shape.element.xml
            import xml.etree.ElementTree as ET
            try:
                root = ET.fromstring(xml_text)

                # 查找图片的描述文本
                for elem in root.iter():
                    if elem.tag.endswith('cNvPr'):  # 图片的非可视属性
                        if elem.get('descr') and elem.get('descr').strip():
                            text_data.append({
                                'slide_index': slide_idx,
                                'shape_index': shape_idx,
                                'text': elem.get('descr').strip(),
                                'shape_obj': shape,
                                'is_image_text': True
                            })
                    elif elem.tag.endswith('title') and elem.text and elem.text.strip():
                        text_data.append({
                            'slide_index': slide_idx,
                            'shape_index': shape_idx,
                            'text': elem.text.strip(),
                            'shape_obj': shape,
                            'is_image_text': True
                        })

            except ET.ParseError:
                pass

    except Exception as e:
        logging.warning(f"[extract_from_media_shape] 媒体形状文本提取失败: {e}")


def _force_cap_font(pptx_path: str,
                    max_title_pt: int = 24,
                    max_body_pt: int = 18):
    """
    把整份 PPT 的字号硬性截到指定上限：
    - 标题占位符 ≤ max_title_pt
    - 其余文本框 ≤ max_body_pt
    """
    prs = Presentation(pptx_path)
    title_types = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            if not shape.has_text_frame:
                continue

            tf = shape.text_frame

            # 1) 判断用哪条上限
            is_title = (
                    shape.is_placeholder
                    and shape.placeholder_format.type in title_types
            )
            limit = max_title_pt if is_title else max_body_pt

            # 2) 禁掉自动放缩，全部按我们写的字号算
            tf.auto_size = MSO_AUTO_SIZE.NONE

            # 3) 段落级 + run 级同时裁剪
            for para in tf.paragraphs:

                # 段落级字体
                if para.font.size is None or para.font.size.pt > limit:
                    para.font.size = Pt(limit)

                # run 级字体
                for run in para.runs:
                    if run.font.size is None or run.font.size.pt > limit:
                        run.font.size = Pt(limit)

    prs.save(pptx_path)
    logging.info(f"[force_cap_font] 字号已硬性限制：标题≤{max_title_pt}pt，正文≤{max_body_pt}pt")


def _repack_pptx(extract_dir: str, output_path: str):
    """重新打包PPTX文件"""
    with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zip_file:
        for root, _, files in os.walk(extract_dir):
            for file in files:
                file_path = os.path.join(root, file)
                arc_name = os.path.relpath(file_path, extract_dir)
                zip_file.write(file_path, arc_name)


def _find_parent_element(root, target_elem):
    """查找元素的父元素"""
    for elem in root.iter():
        for child in elem:
            if child == target_elem:
                return elem
    return None


def _optimize_font_for_english(rPr_elem):
    """优化字体设置以适应英文文本"""
    try:
        # 查找并优化 Latin 字体设置
        latin_elem = None
        ea_elem = None

        for child in rPr_elem:
            if child.tag.endswith('}latin'):
                latin_elem = child
            elif child.tag.endswith('}ea'):
                ea_elem = child

        # 如果存在东亚字体设置，为英文优化 Latin 字体
        if latin_elem is not None:
            # 保持原有的字体设置，但确保使用合适的英文字体
            current_typeface = latin_elem.get('typeface', '微软雅黑')

            # 如果当前是中文字体，确保字体支持英文
            if current_typeface in ['微软雅黑', '宋体', '黑体', '楷体']:
                # 这些字体都支持英文显示，保持原字体设置
                latin_elem.set('typeface', current_typeface)

            # 保持字体大小和其他属性不变

        # 如果有EA字体设置，保持不变（以防混合文本）
        if ea_elem is not None:
            # 保持东亚字体设置不变
            pass

    except Exception as e:
        logging.warning(f"字体优化失败: {e}")


def _adjust_font_for_translation(text_elem, translated_text: str, root):
    """调整翻译后的字体设置，确保格式兼容"""
    try:
        # 更新文本内容
        text_elem.text = translated_text

        # 查找父元素的字体设置 (<a:r> 元素)
        parent = _find_parent_element(root, text_elem)
        if parent is not None and parent.tag.endswith('}r'):
            # 查找字体属性元素 (<a:rPr>)
            rPr_elem = None
            for child in parent:
                if child.tag.endswith('}rPr'):
                    rPr_elem = child
                    break

            if rPr_elem is not None:
                # 调整字体设置以适应英文文本
                _optimize_font_for_english(rPr_elem)
                # 字体大小已经在_adjust_all_font_sizes中调整过了

    except Exception as e:
        logging.warning(f"字体调整失败: {e}")
        # 如果字体调整失败，至少保证文本替换成功
        text_elem.text = translated_text


def _generate_smartart_reference(text_data: List[Dict], translation_dict: Dict[str, str], output_path: str):
    """生成SmartArt翻译参考文件"""
    try:
        smartart_items = [item for item in text_data if item.get('is_smartart_content')]

        if not smartart_items:
            return

        reference_file = output_path.replace('.pptx', '_smartart_translations.txt')

        with open(reference_file, 'w', encoding='utf-8') as f:
            f.write("SmartArt文本翻译参考\n")
            f.write("=" * 50 + "\n\n")
            f.write("注意：SmartArt中的文本内容存储在独立的XML文件中，无法通过程序自动替换。\n")
            f.write("请手动在PowerPoint中修改以下SmartArt文本：\n\n")

            current_slide = -1
            for item in smartart_items:
                if item['slide_index'] != current_slide:
                    current_slide = item['slide_index']
                    f.write(f"幻灯片 {current_slide + 1}:\n")
                    f.write("-" * 20 + "\n")

                original_text = item['text']
                translated_text = translation_dict.get(original_text, original_text)

                if original_text != translated_text:
                    f.write(f"原文: {original_text}\n")
                    f.write(f"译文: {translated_text}\n\n")

        logging.info(f"SmartArt翻译参考文件已生成: {reference_file}")

    except Exception as e:
        logging.warning(f"生成SmartArt参考文件失败: {e}")


def _should_connect_parts(prev_text: str, current_part: str) -> bool:
    """判断两个文本片段是否应该连接"""
    # 如果前一部分明显不完整，应该连接
    incomplete_endings = ['的', '了', '在', '是', '有', '和', '与', '或', '及', '等', '上', '下', '内', '外', '中',
                          '前', '后', '左', '右']
    if any(prev_text.endswith(ending) for ending in incomplete_endings):
        return True

    # 如果当前部分明显是前一部分的续接
    continuation_starts = ['管', '膜', '套', '盒', '袋', '箱', '件', '器', '机', '具', '板', '片', '条', '线', '绳',
                           '带']
    if any(current_part.startswith(start) for start in continuation_starts):
        return True

    # 如果前一部分以动词结尾，当前部分可能是宾语
    verb_endings = ['划', '切', '拿', '放', '拉', '推', '按', '压', '拧', '转', '开', '关', '启', '停']
    if any(prev_text.endswith(verb) for verb in verb_endings):
        return True

    return False


def _smart_combine_text_parts(text_parts: List[str]) -> str:
    """智能组合文本片段"""
    if not text_parts:
        return ""

    if len(text_parts) == 1:
        return _clean_smartart_text(text_parts[0])

    combined = ""
    for i, part in enumerate(text_parts):
        if not part:
            continue

        if i == 0:
            combined = part
        else:
            # 智能连接规则
            # 规则1: 数字、单位直接连接
            if (part.isdigit() or
                    part in ['㎜', '厘米', '米', '克', '公斤', '℃', '°', '%'] or
                    len(part) == 1):
                combined += part

            # 规则2: 标点符号直接连接
            elif part in ['。', '，', '、', '：', '；', '！', '？', ')', '）', '}', '】']:
                combined += part

            # 规则3: 前一个字符是特殊字符，直接连接
            elif (combined.endswith(('＜', '（', '，', '、', '：', '；', '【', '{', '(')) or
                  part.startswith(('缩', '管', '膜', '套'))):  # 常见的连接词
                combined += part

            # 规则4: 句子片段，需要判断是否应该连接
            elif _should_connect_parts(combined, part):
                combined += part

            # 规则5: 独立的完整短语或句子
            else:
                # 如果当前组合的文本已经足够长且完整，就不再添加
                if len(combined) > 8 and ('。' in combined or '，' in combined or len(combined) > 15):
                    break
                combined += part

    return _clean_smartart_text(combined)


def _extract_text_from_diagram_data(xml_data: str, slide_idx: int, shape_idx: int) -> List[Dict]:
    """从图表数据XML中提取文本"""
    texts = []

    try:
        root = ET.fromstring(xml_data)

        # 查找所有文本元素
        text_elements = set()

        # 方法1: 智能组合pt元素中的文本片段
        namespaces = {'dgm': 'http://schemas.openxmlformats.org/drawingml/2006/diagram'}
        pt_elements = root.findall('.//dgm:pt', namespaces)

        for pt in pt_elements:
            # 提取pt元素中的所有文本片段并组合
            t_elements = pt.findall('.//dgm:t', namespaces)
            if t_elements:
                text_parts = []
                for t_elem in t_elements:
                    if t_elem.text and t_elem.text.strip():
                        text_parts.append(t_elem.text.strip())

                if text_parts:
                    # 智能组合文本片段
                    combined_text = _smart_combine_text_parts(text_parts)
                    if combined_text:
                        text_elements.add(combined_text)

        # 方法2: 查找所有独立的文本元素（简化过滤条件）
        for elem in root.iter():
            if elem.text and elem.text.strip():
                text = elem.text.strip()
                if (len(text) > 2 and
                        not text.startswith('{') and
                        text not in ['true', 'false', '0', '1', 'exact', 'lvl'] and
                        not text.isdigit()):
                    # 检查是否包含中文字符或英文字符（用于验证翻译结果）
                    if (any('\u4e00' <= char <= '\u9fff' for char in text) or
                            any(char.isalpha() and ord(char) < 256 for char in text)):
                        text_elements.add(text)

        # 转换为字典格式
        for text in text_elements:
            texts.append({
                'slide_index': slide_idx,
                'shape_index': shape_idx,
                'text': text,
                'is_smartart_content': True,
                'source': 'diagram_data'
            })

    except Exception as e:
        logging.warning(f"图表数据解析错误: {e}")

    return texts


def _extract_smartart_content_from_pptx(pptx_path: str) -> List[Dict]:
    """
    从PPTX文件中提取SmartArt的实际文本内容
    通过直接解析PPTX的内部XML文件来获取SmartArt数据
    """
    smartart_texts = []

    try:
        with zipfile.ZipFile(pptx_path, 'r') as pptx_zip:
            # 获取所有文件列表
            file_list = pptx_zip.namelist()

            # 查找幻灯片文件
            slide_files = [f for f in file_list if f.startswith('ppt/slides/slide') and f.endswith('.xml')]

            for slide_file in slide_files:
                try:
                    # 提取幻灯片编号
                    filename = os.path.basename(slide_file)
                    if filename.startswith('slide') and filename.endswith('.xml'):
                        slide_num_str = filename[5:-4]  # 去掉'slide'和'.xml'
                        slide_num = int(slide_num_str) - 1
                    else:
                        continue
                except (ValueError, IndexError):
                    continue

                # 读取幻灯片XML
                slide_xml = pptx_zip.read(slide_file).decode('utf-8')
                slide_root = ET.fromstring(slide_xml)

                # 查找图形框架（GraphicFrame）
                namespaces = {
                    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                    'dgm': 'http://schemas.openxmlformats.org/drawingml/2006/diagram',
                    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
                }

                # 查找所有图形框架
                graphic_frames = slide_root.findall('.//p:graphicFrame', namespaces)

                for gf_idx, graphic_frame in enumerate(graphic_frames):
                    # 查找关系ID（数据模型关系）
                    rel_ids = []
                    for elem in graphic_frame.iter():
                        for attr, value in elem.attrib.items():
                            if attr.endswith('}dm'):  # 数据模型关系
                                rel_ids.append(value)

                    # 查找对应的关系文件
                    slide_rel_file = slide_file.replace('.xml', '.xml.rels').replace('ppt/slides/',
                                                                                     'ppt/slides/_rels/')

                    if slide_rel_file in file_list and rel_ids:
                        rel_xml = pptx_zip.read(slide_rel_file).decode('utf-8')
                        rel_root = ET.fromstring(rel_xml)

                        # 解析关系文件，找到对应的数据文件
                        for rel_id in rel_ids:
                            relationship = rel_root.find(f".//r:Relationship[@Id='{rel_id}']", {
                                'r': 'http://schemas.openxmlformats.org/package/2006/relationships'})
                            if relationship is not None:
                                target = relationship.get('Target')
                                if target and target.startswith('../diagrams/'):
                                    # 构建完整路径
                                    data_file = f"ppt/diagrams/{target[12:]}"# 去掉'../diagrams/'

                                    if data_file in file_list:
                                        # 读取并解析数据文件
                                        try:
                                            data_xml = pptx_zip.read(data_file).decode('utf-8')
                                            texts = _extract_text_from_diagram_data(data_xml, slide_num,
                                                                                         gf_idx)
                                            smartart_texts.extend(texts)
                                            logging.info(f"从 {data_file} 提取到 {len(texts)} 个文本")
                                        except Exception as e:
                                            logging.warning(f" SmartArt数据文件解析错误: {e}")
                                    else:
                                        logging.warning(f" 数据文件不存在: {data_file}")
                            else:
                                logging.warning(f" 未找到关系ID {rel_id} 对应的关系")
                    else:
                        if not rel_ids:
                            logging.debug(f" 幻灯片{slide_num + 1}无SmartArt数据模型关系")
                        if slide_rel_file not in file_list:
                            logging.debug(f" 关系文件不存在: {slide_rel_file}")

    except Exception as e:
        logging.warning(f" PPTX SmartArt解析错误: {e}")

    return smartart_texts


def _extract_format_info(run):
    """提取文本格式信息 - 增强版本"""
    format_info = {
        'font_name': run.font.name,
        'font_size': run.font.size,
        'font_bold': run.font.bold,
        'font_italic': run.font.italic,
        'font_underline': run.font.underline,
        'font_color': None,
        'color_type': None,
        'color_brightness': None,
        'paragraph_alignment': None
    }
    
    # 提取段落对齐信息
    try:
        if hasattr(run, '_element') and hasattr(run._element, 'getparent'):
            paragraph_element = run._element.getparent()
            if hasattr(paragraph_element, 'getparent'):
                textframe_element = paragraph_element.getparent()
                # 尝试从父段落获取对齐信息
                if hasattr(run, 'paragraph') and hasattr(run.paragraph, 'alignment'):
                    format_info['paragraph_alignment'] = run.paragraph.alignment
    except Exception:
        pass

    # 如果字体大小为None，尝试从段落或父元素获取默认字体大小
    if format_info['font_size'] is None:
        # 尝试从段落获取默认字体大小
        try:
            paragraph = run._parent
            if hasattr(paragraph, 'runs') and paragraph.runs:
                for other_run in paragraph.runs:
                    if other_run.font.size is not None:
                        format_info['font_size'] = other_run.font.size
                        logging.info(f"从其他run获取字体大小: {other_run.font.size}")
                        break
        except:
            pass

        # 如果仍然为None，使用默认值
        if format_info['font_size'] is None:
            from pptx.util import Pt
            # 根据文本类型设置默认字体大小
            default_size = 18  # 默认字体大小
            if run.text and len(run.text.strip()) < 10:  # 短文本可能是标题
                default_size = 24
            format_info['font_size'] = Pt(default_size)
            logging.info(f"设置默认字体大小: {default_size}pt")

    # 安全地获取字体颜色 - 增强版本
    try:
        color_obj = run.font.color
        if color_obj is not None:
            # 尝试获取RGB颜色
            if hasattr(color_obj, 'rgb') and color_obj.rgb is not None:
                format_info['font_color'] = color_obj.rgb
                format_info['color_type'] = 'rgb'
                logging.debug(f"保存RGB颜色: {color_obj.rgb}")
            # 尝试获取主题颜色
            elif hasattr(color_obj, 'theme_color') and color_obj.theme_color is not None:
                format_info['font_color'] = color_obj.theme_color
                format_info['color_type'] = 'theme'
                # 如果有色调信息也保存
                if hasattr(color_obj, 'brightness') and color_obj.brightness is not None:
                    format_info['color_brightness'] = color_obj.brightness
                logging.debug(f"保存主题颜色: {color_obj.theme_color}, 亮度: {format_info.get('color_brightness')}")
    except Exception as e:
        logging.debug(f"获取字体颜色失败: {e}")
        pass

    return format_info







class PPTTranslator:
    def __init__(self, font_size_adjustment: float = 0.6, translator: Translator = SunTranslator()):
        """
        初始化PPT翻译器

        Args:
            font_size_adjustment: 字体大小调整比例（默认0.6，即缩小到60%）
        """
        self.font_size_adjustment = font_size_adjustment
        self.translator = translator

    def _calculate_font_adjustment_ratio(self, original_text: str, translated_text: str) -> float:
        """根据文本长度变化计算字体大小调整比例"""
        if not original_text or not translated_text:
            return self.font_size_adjustment
            
        # 计算字符长度比
        len_ratio = len(translated_text) / len(original_text) if len(original_text) > 0 else 1.0
        
        # 计算视觉宽度比（英文字符通常比中文字符窄）
        original_width = sum(2 if '\u4e00' <= c <= '\u9fff' else 1 for c in original_text)
        translated_width = sum(2 if '\u4e00' <= c <= '\u9fff' else 1 for c in translated_text)
        width_ratio = translated_width / original_width if original_width > 0 else 1.0
        
        # 综合考虑长度和宽度
        combined_ratio = (len_ratio + width_ratio) / 2
        
        # 根据比例计算调整系数
        if combined_ratio <= 1.0:
            # 翻译后更短，不需要缩小
            adjustment = 1.0
        elif combined_ratio <= 1.5:
            # 轻微增长，轻微缩小
            adjustment = 0.85
        elif combined_ratio <= 2.0:
            # 中等增长
            adjustment = 0.7
        elif combined_ratio <= 3.0:
            # 显著增长
            adjustment = 0.6
        else:
            # 极大增长
            adjustment = 0.5
        
        # 确保不会过度缩小
        return max(adjustment, 0.4)

    def _smart_text_wrapping(self, text: str, max_chars_per_line: int = 40) -> str:
        """智能文本换行，优化显示效果"""
        if len(text) <= max_chars_per_line:
            return text
            
        words = text.split()
        lines = []
        current_line = []
        current_length = 0
        
        for word in words:
            word_length = len(word) + (1 if current_line else 0)  # +1 for space
            
            if current_length + word_length > max_chars_per_line and current_line:
                lines.append(' '.join(current_line))
                current_line = [word]
                current_length = len(word)
            else:
                current_line.append(word)
                current_length += word_length
        
        if current_line:
            lines.append(' '.join(current_line))
        
        return '\n'.join(lines)

    def _adjust_textframe_properties(self, text_frame, original_text, translated_text):
        """调整文本框属性以适应翻译后的文本"""
        try:
            from pptx.util import Pt  # 导入Pt
            
            # 计算文本长度变化
            length_ratio = len(translated_text) / len(original_text) if len(original_text) > 0 else 1.0
            
            # 如果文本显著增长，调整文本框属性
            if length_ratio > 1.5:
                # 启用自动调整文本以适应形状
                text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                
                # 调整边距以获得更多空间
                try:
                    text_frame.margin_left = Pt(2)
                    text_frame.margin_right = Pt(2)
                    text_frame.margin_top = Pt(2)
                    text_frame.margin_bottom = Pt(2)
                except:
                    pass
                
                # 启用文本换行
                try:
                    text_frame.word_wrap = True
                except:
                    pass
                    
        except Exception as e:
            logging.warning(f"文本框调整失败: {e}")

    def _is_text_overflow(self, shape) -> bool:
        """检测文本是否溢出形状边界"""
        try:
            if not hasattr(shape, 'text_frame') or not shape.text_frame:
                return False
                
            # 这是一个简化的检测逻辑
            text_length = sum(len(p.text) for p in shape.text_frame.paragraphs)
            
            # 根据形状大小和文本长度估算
            if hasattr(shape, 'width') and hasattr(shape, 'height'):
                shape_area = shape.width * shape.height
                # 粗略估算：每个字符需要的面积
                char_area = 200 * 12700  # 200平方点转换为EMU
                needed_area = text_length * char_area
                
                return needed_area > shape_area * 0.8  # 留20%边距
                
        except:
            return False

    def extract_text_from_ppt(self, ppt_path: str) -> List[Dict]:
        """
        从PPT文件中提取所有文本内容

        Args:
            ppt_path: PPT文件路径

        Returns:
            包含文本信息的列表
        """
        text_data = []

        # 方法1: 使用传统的python-pptx方法
        prs = Presentation(ppt_path)
        for slide_idx, slide in enumerate(prs.slides):
            logging.info(f"处理幻灯片 {slide_idx + 1}/{len(prs.slides)}")

            # 首先检查幻灯片的占位符（标题、内容等）
            try:
                for placeholder in slide.placeholders:
                    if hasattr(placeholder, 'text_frame') and placeholder.text_frame:
                        for para_idx, paragraph in enumerate(placeholder.text_frame.paragraphs):
                            # 处理整个段落的文本
                            if paragraph.text.strip():
                                text_data.append({
                                    'slide_index': slide_idx,
                                    'placeholder_idx': placeholder.placeholder_format.idx,
                                    'paragraph_index': para_idx,
                                    'text': paragraph.text.strip(),
                                    'paragraph_obj': paragraph,
                                    'is_placeholder': True
                                })

                            # 同时处理段落中的每个run（格式化片段）
                            for run_idx, run in enumerate(paragraph.runs):
                                if run.text.strip():
                                    text_data.append({
                                        'slide_index': slide_idx,
                                        'placeholder_idx': placeholder.placeholder_format.idx,
                                        'paragraph_index': para_idx,
                                        'run_index': run_idx,
                                        'text': run.text.strip(),
                                        'run_obj': run,
                                        'paragraph_obj': paragraph,
                                        'is_placeholder': True,
                                        'is_run': True
                                    })
            except Exception as e:
                logging.info(f" 处理占位符时出错: {e}")

            # 使用更全面的形状遍历方法
            self._extract_all_text_from_slide(slide, slide_idx, text_data)

        # 方法2: 提取SmartArt实际内容（新增）
        logging.info("提取SmartArt实际内容...")
        smartart_texts = _extract_smartart_content_from_pptx(ppt_path)
        text_data.extend(smartart_texts)

        return text_data

    def _extract_all_text_from_slide(self, slide, slide_idx: int, text_data: List[Dict]):
        """全面提取幻灯片中的所有文本"""
        for shape_idx, shape in enumerate(slide.shapes):
            try:
                # 递归提取形状中的所有文本
                self._extract_text_from_shape(shape, slide_idx, shape_idx, text_data)

                # 特殊处理：检查形状的所有可能文本属性
                self._extract_text_deep_scan(shape, slide_idx, shape_idx, text_data)

            except Exception as e:
                logging.info(f" 处理形状 {shape_idx} 时出错: {e}")
                continue

    def _extract_text_from_shape(self, shape, slide_idx: int, shape_idx: int, text_data: List[Dict],
                                 parent_info: Dict = None):
        """递归提取形状中的文本"""
        try:
            # 基本文本框检查
            if hasattr(shape, 'text_frame') and shape.text_frame:
                for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    if paragraph.text.strip():
                        item = {
                            'slide_index': slide_idx,
                            'shape_index': shape_idx,
                            'paragraph_index': para_idx,
                            'text': paragraph.text.strip(),
                            'paragraph_obj': paragraph
                        }
                        if parent_info:
                            item.update(parent_info)
                        text_data.append(item)

            # 表格检查 - 增强版，提取所有文本片段
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table = shape.table
                for row_idx, row in enumerate(table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        # 方法1: 提取整个单元格文本（向后兼容）
                        if cell.text.strip():
                            item = {
                                'slide_index': slide_idx,
                                'shape_index': shape_idx,
                                'cell_row': row_idx,
                                'cell_col': col_idx,
                                'text': cell.text.strip(),
                                'cell_obj': cell,
                                'is_cell_full_text': True
                            }
                            if parent_info:
                                item.update(parent_info)
                            text_data.append(item)
                        
                        # 方法2: 详细提取单元格内的每个文本片段
                        if hasattr(cell, 'text_frame') and cell.text_frame:
                            for para_idx, paragraph in enumerate(cell.text_frame.paragraphs):
                                # 提取整个段落文本
                                if paragraph.text.strip():
                                    item = {
                                        'slide_index': slide_idx,
                                        'shape_index': shape_idx,
                                        'cell_row': row_idx,
                                        'cell_col': col_idx,
                                        'paragraph_index': para_idx,
                                        'text': paragraph.text.strip(),
                                        'paragraph_obj': paragraph,
                                        'cell_obj': cell,
                                        'is_table_paragraph': True
                                    }
                                    if parent_info:
                                        item.update(parent_info)
                                    text_data.append(item)
                                
                                # 提取段落内的每个运行（格式化片段）
                                for run_idx, run in enumerate(paragraph.runs):
                                    if run.text.strip():
                                        item = {
                                            'slide_index': slide_idx,
                                            'shape_index': shape_idx,
                                            'cell_row': row_idx,
                                            'cell_col': col_idx,
                                            'paragraph_index': para_idx,
                                            'run_index': run_idx,
                                            'text': run.text.strip(),
                                            'run_obj': run,
                                            'paragraph_obj': paragraph,
                                            'cell_obj': cell,
                                            'is_table_run': True
                                        }
                                        if parent_info:
                                            item.update(parent_info)
                                        text_data.append(item)

            # 直接文本属性检查
            if hasattr(shape, 'text') and shape.text.strip():
                item = {
                    'slide_index': slide_idx,
                    'shape_index': shape_idx,
                    'text': shape.text.strip(),
                    'shape_obj': shape
                }
                if parent_info:
                    item.update(parent_info)
                text_data.append(item)

            # 组合形状递归处理
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub_idx, sub_shape in enumerate(shape.shapes):
                    sub_parent_info = {'is_group': True, 'parent_shape_idx': shape_idx, 'sub_shape_index': sub_idx}
                    if parent_info:
                        sub_parent_info.update(parent_info)
                    self._extract_text_from_shape(sub_shape, slide_idx, shape_idx, text_data, sub_parent_info)

            # SmartArt和图形框架处理
            if hasattr(shape, 'element') and shape.element.tag.endswith('graphicFrame'):
                self._extract_from_smartart_enhanced(text_data, slide_idx, shape_idx, shape)

            # 图片和图表中的文本
            if shape.shape_type in [MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.CHART]:
                _extract_from_media_shape(text_data, slide_idx, shape_idx, shape)

        except Exception as e:
            logging.info(f" 形状文本提取失败: {e}")

    def _extract_text_deep_scan(self, shape, slide_idx: int, shape_idx: int, text_data: List[Dict]):
        """深度扫描形状的所有可能文本位置"""
        try:
            # 检查形状的XML元素，寻找隐藏的文本
            if hasattr(shape, 'element'):
                self._scan_xml_for_text(shape.element, slide_idx, shape_idx, text_data)

            # 检查形状的所有子元素
            if hasattr(shape, '_element'):
                self._scan_xml_for_text(shape._element, slide_idx, shape_idx, text_data)

        except Exception as e:
            logging.info(f" 深度扫描失败: {e}")

    def _scan_xml_for_text(self, element, slide_idx: int, shape_idx: int, text_data: List[Dict]):
        """扫描XML元素寻找文本"""
        try:
            # 递归查找所有包含文本的XML元素
            if element.text and element.text.strip():
                # 过滤掉系统文本和无意义文本
                text = element.text.strip()
                if (len(text) > 1 and
                        not text.startswith('{') and
                        not text.startswith('<') and
                        not text.isdigit() and
                        text not in ['0', '1', 'true', 'false']):
                    text_data.append({
                        'slide_index': slide_idx,
                        'shape_index': shape_idx,
                        'text': text,
                        'xml_element': element,
                        'is_xml_text': True
                    })

            # 递归处理子元素
            for child in element:
                self._scan_xml_for_text(child, slide_idx, shape_idx, text_data)

        except Exception:
            pass  # 忽略XML解析错误

    def _extract_from_smartart_enhanced(self, text_data: List[Dict], slide_idx: int, shape_idx: int, shape):
        """增强版SmartArt文本提取"""
        try:
            # 方法1: 尝试直接获取文本框
            if hasattr(shape, 'text_frame') and shape.text_frame:
                for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    if paragraph.text.strip():
                        text_data.append({
                            'slide_index': slide_idx,
                            'shape_index': shape_idx,
                            'paragraph_index': para_idx,
                            'text': paragraph.text.strip(),
                            'paragraph_obj': paragraph,
                            'is_smartart': True
                        })

            # 方法2: 遍历所有可能的文本运行
            _extract_smartart_runs(shape, slide_idx, shape_idx, text_data)

            # 方法3: XML深度搜索
            if hasattr(shape, 'element'):
                xml_text = shape.element.xml
                import xml.etree.ElementTree as ET
                try:
                    root = ET.fromstring(xml_text)
                    _extract_smartart_xml_text(root, slide_idx, shape_idx, shape, text_data)
                except ET.ParseError:
                    pass

            # 方法4: 尝试从形状的所有子元素中提取
            self._extract_from_shape_children(shape, slide_idx, shape_idx, text_data)

        except Exception as e:
            logging.info(f" 增强SmartArt处理失败: {e}")

    def _extract_from_shape_children(self, shape, slide_idx: int, shape_idx: int, text_data: List[Dict]):
        """从形状的所有子元素中提取文本"""
        try:
            # 尝试访问形状的内部结构
            if hasattr(shape, '_element'):
                element = shape._element
                # 递归搜索所有文本内容
                self._recursive_text_search(element, slide_idx, shape_idx, text_data)
        except Exception as e:
            logging.info(f" 形状子元素文本提取失败: {e}")

    def _recursive_text_search(self, element, slide_idx: int, shape_idx: int, text_data: List[Dict]):
        """递归搜索元素中的所有文本"""
        try:
            if hasattr(element, 'text') and element.text and element.text.strip():
                text = element.text.strip()
                if len(text) > 1 and not text.startswith('{'):
                    text_data.append({
                        'slide_index': slide_idx,
                        'shape_index': shape_idx,
                        'text': text,
                        'element_obj': element,
                        'is_recursive_text': True
                    })

            # 递归处理子元素
            if hasattr(element, '__iter__'):
                for child in element:
                    self._recursive_text_search(child, slide_idx, shape_idx, text_data)
        except Exception:
            pass

    def _extract_from_group_shape(self, text_data: List[Dict], slide_idx: int, shape_idx: int, group_shape):
        """从组合形状中提取文本"""
        try:
            for sub_shape_idx, sub_shape in enumerate(group_shape.shapes):
                if hasattr(sub_shape, 'text_frame') and sub_shape.text_frame:
                    for para_idx, paragraph in enumerate(sub_shape.text_frame.paragraphs):
                        if paragraph.text.strip():
                            text_data.append({
                                'slide_index': slide_idx,
                                'shape_index': shape_idx,
                                'sub_shape_index': sub_shape_idx,
                                'paragraph_index': para_idx,
                                'text': paragraph.text.strip(),
                                'paragraph_obj': paragraph,
                                'is_group': True
                            })
                # 递归处理嵌套的组合形状
                elif sub_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    self._extract_from_group_shape(text_data, slide_idx, shape_idx, sub_shape)
        except Exception as e:
            logging.info(f" 处理组合形状时出错: {e}")

    def _extract_from_smartart(self, text_data: List[Dict], slide_idx: int, shape_idx: int, shape):
        """从SmartArt图形中提取文本"""
        try:
            # SmartArt可能包含在graphicFrame中
            if hasattr(shape, 'element'):
                # 尝试从XML中提取文本
                xml_text = shape.element.xml

                # 查找所有文本内容
                import xml.etree.ElementTree as ET
                try:
                    root = ET.fromstring(xml_text)
                    # 查找所有包含文本的元素
                    for elem in root.iter():
                        if elem.text and elem.text.strip():
                            # 过滤掉一些非用户文本
                            if not elem.text.strip().startswith('{') and len(elem.text.strip()) > 1:
                                text_data.append({
                                    'slide_index': slide_idx,
                                    'shape_index': shape_idx,
                                    'text': elem.text.strip(),
                                    'smartart_obj': shape,
                                    'xml_element': elem,
                                    'is_smartart': True
                                })
                except ET.ParseError:
                    pass

            # 另一种方法：尝试获取SmartArt的文本框
            if hasattr(shape, 'text_frame') and shape.text_frame:
                for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    if paragraph.text.strip():
                        text_data.append({
                            'slide_index': slide_idx,
                            'shape_index': shape_idx,
                            'paragraph_index': para_idx,
                            'text': paragraph.text.strip(),
                            'paragraph_obj': paragraph,
                            'is_smartart': True
                        })
        except Exception as e:
            logging.warning(f"处理SmartArt时出错: {e}")

    def _apply_format(self, run, format_info):
        """应用文本格式（使用固定比例）"""
        try:
            # 字体名称
            if format_info.get('font_name'):
                run.font.name = format_info['font_name']

            # 字体大小 - 关键修复点，应用字体大小调整比例
            if format_info.get('font_size') is not None:
                orig = format_info['font_size']
                # 支持 Length 对象 / int EMU / float
                original_pt = orig.pt if hasattr(orig, 'pt') else (orig/12700 if isinstance(orig, int) else float(orig))
                # 只缩不放
                adjusted_pt = int(original_pt * self.font_size_adjustment)
                if adjusted_pt > original_pt:
                    adjusted_pt = int(original_pt)
                run.font.size = Pt(max(1, adjusted_pt))
                logging.debug(f"字体大小调整: {original_pt}pt → {adjusted_pt}pt")
            else:
                # run.font.size=None，跳过让 PPT 自己处理
                logging.debug("无原始字体大小，跳过设置")
                
            # 其他格式属性
            if format_info.get('font_bold') is not None:
                run.font.bold = format_info['font_bold']
            if format_info.get('font_italic') is not None:
                run.font.italic = format_info['font_italic']
            if format_info.get('font_underline') is not None:
                run.font.underline = format_info['font_underline']
                
        except Exception as e:
            logging.warning(f"格式应用失败: {e}")

    def _smart_alignment_optimization(self, paragraph, original_text, translated_text, original_alignment):
        """智能对齐优化"""
        try:
            from pptx.enum.text import PP_ALIGN
            
            # 如果原来是左对齐，且翻译后文本明显更长，考虑使用两端对齐或居中
            if original_alignment == PP_ALIGN.LEFT:
                length_ratio = len(translated_text) / len(original_text) if len(original_text) > 0 else 1.0
                
                # 如果翻译后文本显著更长（通常中译英的情况）
                if length_ratio > 1.5:
                    # 对于较长的文本（超过50个字符），使用两端对齐
                    if len(translated_text) > 50:
                        paragraph.alignment = PP_ALIGN.JUSTIFY
                        logging.info(f"对齐优化: 长文本使用两端对齐")
                    # 对于中等长度文本，使用居中对齐
                    elif len(translated_text) > 20:
                        paragraph.alignment = PP_ALIGN.CENTER
                        logging.info(f"对齐优化: 中等文本使用居中对齐")
                    else:
                        # 短文本保持原对齐
                        paragraph.alignment = original_alignment
                        
                # 如果翻译后文本变短（通常英译中的情况），保持原对齐
                elif length_ratio < 0.8:
                    paragraph.alignment = original_alignment
                    
                # 长度相近，使用智能判断
                else:
                    # 如果文本较长，建议两端对齐
                    if len(translated_text) > 80:
                        paragraph.alignment = PP_ALIGN.JUSTIFY
                        logging.info(f"对齐优化: 长文本使用两端对齐")
                    else:
                        paragraph.alignment = original_alignment
                        
            # 右对齐和居中对齐通常保持不变
            elif original_alignment in [PP_ALIGN.RIGHT, PP_ALIGN.CENTER]:
                paragraph.alignment = original_alignment
                
            # 两端对齐保持不变
            elif original_alignment == PP_ALIGN.JUSTIFY:
                paragraph.alignment = original_alignment
                
            else:
                # 默认情况，保持原对齐
                paragraph.alignment = original_alignment
                
        except Exception as e:
            logging.debug(f"对齐优化失败: {e}")
            # 失败时保持原对齐
            if original_alignment is not None:
                paragraph.alignment = original_alignment

    def _apply_format_with_smart_sizing(self, run, format_info, original_text, translated_text):
        """应用格式时智能调整字体大小和对齐方式"""
        try:
            from pptx.util import Pt  # 导入Pt
            
            # 计算动态调整比例
            adjustment_ratio = self._calculate_font_adjustment_ratio(original_text, translated_text)
            
            # 字体名称
            if format_info.get('font_name'):
                run.font.name = format_info['font_name']
            
            # 字体大小 - 使用动态调整
            if format_info.get('font_size') is not None:
                orig = format_info['font_size']
                original_pt = orig.pt if hasattr(orig, 'pt') else (orig/12700 if isinstance(orig, int) else float(orig))
                
                # 应用动态调整比例（只缩不放）
                adjusted_pt = min(int(original_pt * adjustment_ratio), int(original_pt))
                
                # 设置最小字体大小限制
                min_font_size = 8
                adjusted_pt = max(min_font_size, adjusted_pt)
                
                run.font.size = Pt(adjusted_pt)
                logging.debug(f"动态字体调整: {original_pt}pt → {adjusted_pt}pt (比例: {adjustment_ratio:.2f})")
            else:
                logging.debug("无原始字体大小，跳过动态调整")

            # 其他格式属性
            if format_info.get('font_bold') is not None:
                run.font.bold = format_info['font_bold']
            if format_info.get('font_italic') is not None:
                run.font.italic = format_info['font_italic']
            if format_info.get('font_underline') is not None:
                run.font.underline = format_info['font_underline']
                
            # 字体颜色
            if format_info.get('font_color'):
                try:
                    color_type = format_info.get('color_type', 'rgb')
                    if color_type == 'rgb':
                        run.font.color.rgb = format_info['font_color']
                        logging.debug(f"应用RGB颜色: {format_info['font_color']}")
                    elif color_type == 'theme':
                        run.font.color.theme_color = format_info['font_color']
                        if format_info.get('brightness'):
                            run.font.color.brightness = format_info['brightness']
                        logging.debug(f"应用主题颜色: {format_info['font_color']}")
                except Exception as e:
                    logging.warning(f"颜色设置失败: {e}")
            
            # 智能对齐优化
            try:
                if hasattr(run, 'paragraph') and format_info.get('paragraph_alignment') is not None:
                    self._smart_alignment_optimization(
                        run.paragraph, 
                        original_text, 
                        translated_text, 
                        format_info['paragraph_alignment']
                    )
            except Exception as e:
                logging.debug(f"对齐优化失败: {e}")
                    
        except Exception as e:
            logging.warning(f"智能格式应用失败: {e}")

            # 字体颜色 - 增强版本
            if format_info.get('font_color'):
                try:
                    color_type = format_info.get('color_type', 'rgb')
                    if color_type == 'rgb':
                        run.font.color.rgb = format_info['font_color']
                        logging.debug(f"应用RGB颜色: {format_info['font_color']}")
                    elif color_type == 'theme':
                        run.font.color.theme_color = format_info['font_color']
                        # 应用亮度调整
                        if format_info.get('color_brightness') is not None:
                            run.font.color.brightness = format_info['color_brightness']
                        logging.debug(f"应用主题颜色: {format_info['font_color']}, 亮度: {format_info.get('color_brightness')}")
                except Exception as e:
                    logging.debug(f"应用字体颜色失败: {e}")
                    pass

        except Exception as e:
            logging.info(f" 格式应用失败: {e}")
            # 如果其他属性失败，至少确保字体大小正确
            try:
                if format_info.get('font_size') is not None:
                    from pptx.util import Pt
                    original_size = format_info['font_size']
                    if hasattr(original_size, 'pt'):
                        original_pt = original_size.pt
                    else:
                        original_pt = float(original_size) if not isinstance(original_size,
                                                                             int) else original_size / 12700
                    adjusted_size_value = max(8, min(int(original_pt * self.font_size_adjustment), int(original_pt)))
                    run.font.size = Pt(adjusted_size_value)
            except:
                pass
    
    def _apply_smart_alignment_to_paragraph(self, paragraph, original_text, translated_text):
        """为段落应用智能对齐优化"""
        try:
            from pptx.enum.text import PP_ALIGN
            # 获取当前对齐方式
            original_alignment = paragraph.alignment
            self._smart_alignment_optimization(paragraph, original_text, translated_text, original_alignment)
        except Exception as e:
            logging.debug(f"段落对齐优化失败: {e}")
    
    def _optimize_textframe_for_wrapping(self, text_frame, original_text, translated_text):
        """优化文本框以支持自动换行和填充"""
        try:
            from pptx.enum.text import MSO_AUTO_SIZE
            from pptx.util import Pt
            
            # 计算文本长度比例
            length_ratio = len(translated_text) / len(original_text) if len(original_text) > 0 else 1.0
            
            # 启用自动换行
            try:
                text_frame.word_wrap = True
                logging.debug("启用文本框自动换行")
            except Exception as e:
                logging.debug(f"设置自动换行失败: {e}")
            
            # 根据文本长度变化调整文本框设置
            if length_ratio > 1.2:  # 文本变长
                # 设置自动调整以适应文本
                try:
                    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    logging.debug("设置文本框自动调整大小")
                except Exception as e:
                    logging.debug(f"设置自动调整失败: {e}")
                
                # 减少边距以获得更多文本空间
                try:
                    text_frame.margin_left = Pt(4)
                    text_frame.margin_right = Pt(4)
                    text_frame.margin_top = Pt(2)
                    text_frame.margin_bottom = Pt(2)
                    logging.debug("调整文本框边距")
                except Exception as e:
                    logging.debug(f"设置文本框边距失败: {e}")
                    
            elif length_ratio < 0.8:  # 文本变短
                # 对于变短的文本，保持原始设置但启用换行
                try:
                    text_frame.auto_size = MSO_AUTO_SIZE.NONE
                except Exception as e:
                    logging.debug(f"设置固定大小失败: {e}")
            
            else:  # 长度相近
                # 保持合理的自动调整
                try:
                    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                except Exception as e:
                    logging.debug(f"设置文本适应形状失败: {e}")
                    
        except Exception as e:
            logging.debug(f"文本框优化失败: {e}")
    
    def _set_paragraph_line_spacing(self, paragraph, line_spacing_factor=1.0):
        """设置段落行间距以优化文本显示"""
        try:
            from pptx.util import Pt
            
            # 设置行间距为单倍行距或稍微紧密一些
            if hasattr(paragraph, 'line_spacing'):
                paragraph.line_spacing = line_spacing_factor
                logging.debug(f"设置段落行间距: {line_spacing_factor}")
                
            # 减少段落间距
            if hasattr(paragraph, 'space_before'):
                paragraph.space_before = Pt(0)
            if hasattr(paragraph, 'space_after'):
                paragraph.space_after = Pt(2)
                
        except Exception as e:
            logging.debug(f"设置段落间距失败: {e}")

    def replace_text_in_ppt(self, ppt_path: str, output_path: str,
                            text_data: List[Dict], translation_dict: Dict[str, str]):
        """
        在PPT中替换文本，保持格式
        """
        prs = Presentation(ppt_path)
        replaced_count = 0

        # 去重处理，避免重复替换
        processed_items = set()

        for item in text_data:
            original_text = item['text']
            translated_text = translation_dict.get(original_text, original_text)

            if original_text == translated_text:
                continue

            # 创建唯一标识符避免重复处理
            item_key = f"{item['slide_index']}_{item.get('shape_index', 'no_shape')}_{item.get('paragraph_index', 'x')}"
            if item_key in processed_items:
                continue
            processed_items.add(item_key)

            try:
                # 处理占位符文本
                if item.get('is_placeholder'):
                    slide = prs.slides[item['slide_index']]
                    placeholder = slide.placeholders[item['placeholder_idx']]
                    paragraph = placeholder.text_frame.paragraphs[item['paragraph_index']]
                    translated = translation_dict[item['text']]
                    # 清空整个 text_frame（所有段落和 run）
                    placeholder.text_frame.clear()  
                    # 直接赋值，会自动产生一个 run
                    placeholder.text = translated  
                    # 处理单个run（格式化片段）
                    if item.get('is_run') and 'run_obj' in item:
                        run = item['run_obj']
                        # 保存原始格式信息
                        format_info = _extract_format_info(run)
                        # 替换文本
                        run.text = translated_text
                        # 应用格式（包括调整后的字体大小）
                        self._apply_format(run, format_info)
                        logging.info(f" 占位符run替换: '{original_text}' → '{translated_text}'")

                    # 处理整个段落
                    elif not item.get('is_run'):
                        if paragraph.runs:
                            # 收集所有run的格式信息
                            runs_format = []
                            for run in paragraph.runs:
                                runs_format.append(_extract_format_info(run))

                            # 如果翻译后的文本长度与原文本相近，保持多run结构
                            if len(paragraph.runs) > 1 and len(translated_text) > len(original_text) * 0.8:
                                # 尝试保持多run结构并调整所有字体大小
                                for i, run in enumerate(paragraph.runs):
                                    if i < len(runs_format):
                                        self._apply_format(run, runs_format[i])
                                # 只更新第一个run的文本，其他run清空
                                paragraph.runs[0].text = translated_text
                                for i in range(1, len(paragraph.runs)):
                                    paragraph.runs[i].text = ""
                            else:
                                # 使用第一个run的格式重建整个段落
                                main_format = runs_format[0] if runs_format else {}
                                paragraph.clear()
                                run = paragraph.add_run()
                                run.text = translated_text
                                self._apply_format(run, main_format)

                            logging.info(f" 占位符段落替换: '{original_text}' → '{translated_text}'")
                        else:
                            new_run = paragraph.add_run()
                            new_run.text = translated_text
                            # 应用默认字体大小调整
                            try:
                                from pptx.util import Pt
                                new_run.font.size = Pt(min(int(18 * self.font_size_adjustment), 18))  # 默认18pt，只缩不放
                            except:
                                pass

                # 处理普通文本框
                elif 'paragraph_obj' in item and not any(
                        item.get(key) for key in
                        ['is_group', 'is_smartart', 'is_special', 'is_xml_text', 'is_image_text']):
                    # 重新定位到正确的段落对象
                    try:
                        if 'shape_index' not in item:
                            logging.info(f" 缺少shape_index: {original_text}")
                            continue
                        slide = prs.slides[item['slide_index']]
                        shape = slide.shapes[item['shape_index']]
                        
                        # 检查形状是否有text_frame属性
                        if not hasattr(shape, 'text_frame') or not shape.text_frame:
                            logging.warning(f" 形状无text_frame属性: {original_text} - {type(shape).__name__}")
                            continue
                            
                        paragraph = shape.text_frame.paragraphs[item['paragraph_index']]
                    except (IndexError, AttributeError) as e:
                        logging.info(f" 无法定位段落对象: {original_text} - {e}")
                        continue

                    if paragraph.runs:
                        # 保存所有run的格式信息
                        runs_format = []
                        for run in paragraph.runs:
                            format_info = {
                                'font_name': run.font.name,
                                'font_size': run.font.size,
                                'font_bold': run.font.bold,
                                'font_italic': run.font.italic,
                                'font_underline': run.font.underline,
                                'font_color': None
                            }

                            # 安全地获取字体颜色 - 增强版本
                            try:
                                color_obj = run.font.color
                                if color_obj is not None:
                                    # 尝试获取RGB颜色
                                    if hasattr(color_obj, 'rgb') and color_obj.rgb is not None:
                                        format_info['font_color'] = color_obj.rgb
                                        format_info['color_type'] = 'rgb'
                                    # 尝试获取主题颜色
                                    elif hasattr(color_obj, 'theme_color') and color_obj.theme_color is not None:
                                        format_info['font_color'] = color_obj.theme_color
                                        format_info['color_type'] = 'theme'
                                        # 如果有色调信息也保存
                                        if hasattr(color_obj, 'brightness') and color_obj.brightness is not None:
                                            format_info['color_brightness'] = color_obj.brightness
                            except Exception as e:
                                logging.debug(f"获取字体颜色失败: {e}")
                                pass

                            runs_format.append(format_info)

                        # 使用第一个run的格式作为主格式
                        main_format = runs_format[0] if runs_format else {}

                        # 清空并重新设置
                        paragraph.clear()
                        run = paragraph.add_run()
                        run.text = translated_text

                        # 恢复格式，特别注意字体大小的精确保持
                        try:
                            # 字体名称
                            if main_format.get('font_name'):
                                run.font.name = main_format['font_name']

                            # 字体大小 - 关键修复点，应用字体大小调整比例
                            if main_format.get('font_size') is not None:
                                original_size = main_format['font_size']

                                # 处理不同类型的字体大小值并应用调整比例
                                if hasattr(original_size, 'pt'):
                                    # 如果是pptx Length对象
                                    original_pt = original_size.pt
                                elif isinstance(original_size, int):
                                    # 如果是整数值（EMU单位），转换为点
                                    original_pt = original_size / 12700  # 1 pt = 12700 EMU
                                else:
                                    # 其他情况，尝试直接使用
                                    original_pt = float(original_size)

                                # 应用字体大小调整比例（只缩不放）
                                adjusted_size_value = min(int(original_pt * self.font_size_adjustment), int(original_pt))
                                from pptx.util import Pt
                                adjusted_size = Pt(adjusted_size_value)
                                run.font.size = adjusted_size
                                logging.info(
                                    f"字体大小调整: {original_pt}pt → {adjusted_size_value}pt (比例: {self.font_size_adjustment})")
                                # 验证设置是否成功
                                if run.font.size != adjusted_size:
                                    logging.warning(f" 字体大小设置异常: 期望 {adjusted_size}, 实际 {run.font.size}")

                            # 其他格式属性
                            if main_format.get('font_bold') is not None:
                                run.font.bold = main_format['font_bold']
                            if main_format.get('font_italic') is not None:
                                run.font.italic = main_format['font_italic']
                            if main_format.get('font_underline') is not None:
                                run.font.underline = main_format['font_underline']

                            # 字体颜色 - 增强版本
                            if main_format.get('font_color'):
                                try:
                                    color_type = main_format.get('color_type', 'rgb')
                                    if color_type == 'rgb':
                                        run.font.color.rgb = main_format['font_color']
                                    elif color_type == 'theme':
                                        run.font.color.theme_color = main_format['font_color']
                                        # 应用亮度调整
                                        if main_format.get('color_brightness') is not None:
                                            run.font.color.brightness = main_format['color_brightness']
                                except Exception as e:
                                    logging.debug(f"应用字体颜色失败: {e}")
                                    pass

                        except Exception as e:
                            logging.warning(f" 格式设置失败: {e}")
                            # 如果其他属性失败，至少确保字体大小正确
                            try:
                                if main_format.get('font_size') is not None:
                                    run.font.size = main_format['font_size']
                            except:
                                pass
                    else:
                        paragraph.add_run().text = translated_text

                # 处理表格单元格
                elif 'cell_obj' in item:
                    # 重新定位到正确的单元格对象
                    try:
                        if 'shape_index' not in item:
                            logging.warning(f" 表格缺少shape_index: {original_text}")
                            continue
                        slide = prs.slides[item['slide_index']]
                        shape = slide.shapes[item['shape_index']]
                        table = shape.table
                        cell = table.rows[item['cell_row']].cells[item['cell_col']]
                        cell.text = translated_text
                    except (IndexError, AttributeError) as e:
                        logging.warning(f" 表格文本替换失败: {original_text} - {e}")
                        continue

                # 处理图形文本
                elif 'shape_obj' in item:
                    try:
                        if 'shape_index' not in item:
                            logging.warning(f" 图形缺少shape_index: {original_text}")
                            continue
                        slide = prs.slides[item['slide_index']]
                        shape = slide.shapes[item['shape_index']]
                        if hasattr(shape, 'text'):
                            shape.text = translated_text
                    except (IndexError, AttributeError) as e:
                        logging.warning(f" 图形文本替换失败: {original_text} - {e}")
                        continue

                # 处理组合形状中的文本
                elif item.get('is_group'):
                    try:
                        if 'shape_index' not in item:
                            logging.warning(f" 组合形状缺少shape_index: {original_text}")
                            continue
                        slide = prs.slides[item['slide_index']]
                        group_shape = slide.shapes[item['shape_index']]
                        sub_shape = group_shape.shapes[item['sub_shape_index']]
                        # 检查子形状是否有text_frame属性
                        if not hasattr(sub_shape, 'text_frame') or sub_shape.text_frame is None:
                            logging.warning(f" 组合形状子形状无text_frame: {original_text}")
                            continue
                        paragraph = sub_shape.text_frame.paragraphs[item['paragraph_index']]
                    except (IndexError, AttributeError) as e:
                        logging.warning(f" 组合形状文本替换失败: {original_text} - {e}")
                        continue

                    if paragraph.runs:
                        # 保存格式信息
                        main_format = _extract_format_info(paragraph.runs[0])

                        # 替换文本
                        paragraph.clear()
                        run = paragraph.add_run()
                        run.text = translated_text

                        # 恢复格式
                        self._apply_format(run, main_format)

                # 处理SmartArt文本
                elif item.get('is_smartart'):
                    try:
                        # 方法1: 处理文本运行
                        if item.get('is_run') and 'run_obj' in item:
                            run = item['run_obj']
                            format_info = _extract_format_info(run)
                            run.text = translated_text
                            self._apply_format(run, format_info)

                        # 方法2: 处理段落对象
                        elif 'paragraph_obj' in item:
                            paragraph = item['paragraph_obj']
                            if paragraph.runs:
                                main_format = _extract_format_info(paragraph.runs[0])
                                paragraph.clear()
                                run = paragraph.add_run()
                                run.text = translated_text
                                self._apply_format(run, main_format)
                            else:
                                paragraph.add_run().text = translated_text

                        # 方法3: 处理XML元素
                        elif 'xml_element' in item:
                            try:
                                item['xml_element'].text = translated_text
                            except:
                                logging.error(f" SmartArt XML文本替换失败: {original_text}")

                        # 方法4: 处理元素对象
                        elif 'element_obj' in item:
                            try:
                                item['element_obj'].text = translated_text
                            except:
                                logging.error(f" SmartArt元素文本替换失败: {original_text}")

                        # 方法5: 备用方法 - 尝试重新定位并替换
                        else:
                            self._fallback_smartart_replace(prs, item, translated_text, original_text)

                    except Exception as e:
                        logging.info(f" SmartArt文本替换失败: {original_text} - {e}")

                # 处理特殊形状
                elif item.get('is_special'):
                    if 'paragraph_obj' in item:
                        paragraph = item['paragraph_obj']
                        if paragraph.runs:
                            main_format = _extract_format_info(paragraph.runs[0])
                            paragraph.clear()
                            run = paragraph.add_run()
                            run.text = translated_text
                            self._apply_format(run, main_format)
                        else:
                            paragraph.add_run().text = translated_text
                    elif 'shape_obj' in item:
                        try:
                            if 'shape_index' not in item:
                                logging.info(f" 特殊形状缺少shape_index: {original_text}")
                                continue
                            slide = prs.slides[item['slide_index']]
                            shape = slide.shapes[item['shape_index']]
                            if hasattr(shape, 'text'):
                                shape.text = translated_text
                        except (IndexError, AttributeError) as e:
                            logging.info(f" 特殊形状文本替换失败: {original_text} - {e}")
                            continue

                # 处理XML文本
                elif item.get('is_xml_text'):
                    try:
                        if 'xml_element' in item:
                            # 更新文本内容
                            item['xml_element'].text = translated_text

                            # 尝试调整XML中的字体大小
                            self._adjust_xml_font_size(item['xml_element'])

                            replaced_count += 1
                            logging.info(f" XML文本替换: '{original_text}' → '{translated_text}'")
                    except Exception as e:
                        logging.info(f" XML文本替换失败: {original_text} - {e}")

                # 处理图片文本
                elif item.get('is_image_text'):
                    try:
                        # 图片的描述文本通常无法直接修改，记录但跳过
                        logging.info(f"ℹ️ 图片描述文本: '{original_text}' → '{translated_text}' (仅记录)")
                    except:
                        logging.info(f" 图片文本处理失败: {original_text}")

                # 处理SmartArt内容文本
                elif item.get('is_smartart_content'):
                    # SmartArt内容文本存储在独立的XML文件中，无法直接修改
                    # 这里记录翻译信息，提醒用户手动处理
                    logging.info(f"ℹ️ SmartArt内容文本: '{original_text}' → '{translated_text}' (需要手动修改)")
                    # 可以考虑将这些信息保存到文件中供用户参考
                    replaced_count += 1

                # 处理组合形状中的文本（增强版）
                elif item.get('is_group') and item.get('parent_shape_idx') is not None:
                    try:
                        slide = prs.slides[item['slide_index']]
                        parent_shape = slide.shapes[item['parent_shape_idx']]
                        if hasattr(parent_shape, 'shapes') and item.get('sub_shape_index') is not None:
                            sub_shape = parent_shape.shapes[item['sub_shape_index']]
                            if 'paragraph_obj' in item:
                                paragraph = item['paragraph_obj']
                                if paragraph.runs:
                                    main_format = _extract_format_info(paragraph.runs[0])
                                    paragraph.clear()
                                    run = paragraph.add_run()
                                    run.text = translated_text
                                    self._apply_format(run, main_format)
                    except Exception as e:
                        logging.info(f" 增强组合形状文本替换失败: {original_text} - {e}")

                replaced_count += 1
                logging.info(f" 替换: '{original_text}' → '{translated_text}'")

            except Exception as e:
                logging.info(f" 替换失败: {original_text} - {str(e)}")
                continue

        # 如果替换数量太少，使用备用的简单替换策略
        if replaced_count < len(
                [item for item in text_data if translation_dict.get(item['text'], item['text']) != item['text']]) * 0.3:
            logging.info(f" 替换数量偏少 ({replaced_count})，使用备用替换策略...")
            replaced_count += self._simple_text_replacement(prs, translation_dict)

        prs.save(output_path)
        logging.info(f" PPT翻译完成！共替换 {replaced_count} 处文本")

    def _simple_text_replacement(self, prs, translation_dict: Dict[str, str]) -> int:
        """简单的文本替换策略（备用方案）"""
        replaced_count = 0

        for slide in prs.slides:
            # 首先处理占位符（标题、内容等）
            try:
                for placeholder in slide.placeholders:
                    if hasattr(placeholder, 'text_frame') and placeholder.text_frame:
                        for paragraph in placeholder.text_frame.paragraphs:
                            # 处理段落中的每个run
                            for run in paragraph.runs:
                                if run.text.strip():
                                    original_text = run.text.strip()
                                    if original_text in translation_dict:
                                        translated_text = translation_dict[original_text]
                                        if original_text != translated_text:
                                            try:
                                                format_info = _extract_format_info(run)
                                                run.text = translated_text
                                                self._apply_format(run, format_info)
                                                replaced_count += 1
                                                logging.info(f"备用占位符run替换: '{original_text}' → '{translated_text}'")
                                            except Exception as e:
                                                logging.info(f"备用占位符替换失败: {original_text} - {e}")

                            # 处理整个段落
                            if paragraph.text.strip():
                                original_text = paragraph.text.strip()
                                if original_text in translation_dict:
                                    translated_text = translation_dict[original_text]
                                    if original_text != translated_text:
                                        try:
                                            if paragraph.runs:
                                                first_run = paragraph.runs[0]
                                                format_info = _extract_format_info(first_run)
                                                paragraph.clear()
                                                new_run = paragraph.add_run()
                                                new_run.text = translated_text
                                                self._apply_format(new_run, format_info)
                                                replaced_count += 1
                                                logging.info(f"🔄 备用占位符段落替换: '{original_text}' → '{translated_text}'")
                                        except Exception as e:
                                            logging.info(f" 备用占位符段落替换失败: {original_text} - {e}")
            except Exception as e:
                logging.info(f" 备用占位符处理失败: {e}")

            # 处理其他形状
            for shape in slide.shapes:
                try:
                    # 处理文本框
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            # 先处理段落中的每个run
                            for run in paragraph.runs:
                                if run.text.strip():
                                    original_text = run.text.strip()
                                    if original_text in translation_dict:
                                        translated_text = translation_dict[original_text]
                                        if original_text != translated_text:
                                            try:
                                                format_info = _extract_format_info(run)
                                                run.text = translated_text
                                                self._apply_format(run, format_info)
                                                replaced_count += 1
                                                logging.info(f"🔄 备用run替换: '{original_text}' → '{translated_text}'")
                                            except Exception as e:
                                                logging.info(f" 备用run替换失败: {original_text} - {e}")

                            # 再处理整个段落
                            if paragraph.text.strip():
                                original_text = paragraph.text.strip()
                                if original_text in translation_dict:
                                    translated_text = translation_dict[original_text]
                                    if original_text != translated_text:
                                        try:
                                            # 保存第一个run的格式
                                            if paragraph.runs:
                                                first_run = paragraph.runs[0]
                                                format_info = _extract_format_info(first_run)

                                                # 替换文本
                                                paragraph.clear()
                                                new_run = paragraph.add_run()
                                                new_run.text = translated_text
                                                self._apply_format(new_run, format_info)
                                                replaced_count += 1
                                                logging.info(f"🔄 备用段落替换: '{original_text}' → '{translated_text}'")
                                        except Exception as e:
                                            logging.info(f" 备用段落替换失败: {original_text} - {e}")

                    # 处理表格
                    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        table = shape.table
                        for row in table.rows:
                            for cell in row.cells:
                                # 处理单元格中的每个run
                                if hasattr(cell, 'text_frame') and cell.text_frame:
                                    for paragraph in cell.text_frame.paragraphs:
                                        for run in paragraph.runs:
                                            if run.text.strip():
                                                original_text = run.text.strip()
                                                if original_text in translation_dict:
                                                    translated_text = translation_dict[original_text]
                                                    if original_text != translated_text:
                                                        try:
                                                            format_info = _extract_format_info(run)
                                                            run.text = translated_text
                                                            self._apply_format(run, format_info)
                                                            replaced_count += 1
                                                            logging.info(
                                                                f"🔄 备用表格run替换: '{original_text}' → '{translated_text}'")
                                                        except Exception as e:
                                                            logging.info(f" 备用表格run替换失败: {original_text} - {e}")

                                # 处理整个单元格
                                if cell.text.strip():
                                    original_text = cell.text.strip()
                                    if original_text in translation_dict:
                                        translated_text = translation_dict[original_text]
                                        if original_text != translated_text:
                                            try:
                                                # 保存原始格式并应用调整
                                                if hasattr(cell,
                                                           'text_frame') and cell.text_frame and cell.text_frame.paragraphs:
                                                    first_paragraph = cell.text_frame.paragraphs[0]
                                                    if first_paragraph.runs:
                                                        format_info = _extract_format_info(first_paragraph.runs[0])
                                                        cell.text = translated_text
                                                        # 重新应用格式
                                                        if first_paragraph.runs:
                                                            self._apply_format(first_paragraph.runs[0], format_info)
                                                    else:
                                                        cell.text = translated_text
                                                else:
                                                    cell.text = translated_text
                                                replaced_count += 1
                                                logging.info(f"🔄 备用表格替换: '{original_text}' → '{translated_text}'")
                                            except Exception as e:
                                                logging.info(f" 备用表格替换失败: {original_text} - {e}")

                    # 处理组合形状
                    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        replaced_count += self._replace_text_in_group_enhanced(shape, translation_dict)

                except Exception as e:
                    logging.info(f" 备用处理形状失败: {e}")
                    continue

        return replaced_count

    def _replace_text_in_group(self, group_shape, translation_dict: Dict[str, str]) -> int:
        """在组合形状中替换文本"""
        replaced_count = 0

        try:
            for sub_shape in group_shape.shapes:
                if hasattr(sub_shape, 'text_frame') and sub_shape.text_frame:
                    for paragraph in sub_shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            original_text = paragraph.text.strip()
                            if original_text in translation_dict:
                                translated_text = translation_dict[original_text]
                                if original_text != translated_text:
                                    try:
                                        if paragraph.runs:
                                            first_run = paragraph.runs[0]
                                            format_info = _extract_format_info(first_run)

                                            paragraph.clear()
                                            new_run = paragraph.add_run()
                                            new_run.text = translated_text
                                            self._apply_format(new_run, format_info)
                                            replaced_count += 1
                                            logging.info(f"🔄 备用组合替换: '{original_text}' → '{translated_text}'")
                                    except Exception as e:
                                        logging.info(f" 备用组合替换失败: {original_text} - {e}")

                # 递归处理嵌套组合
                elif sub_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    replaced_count += self._replace_text_in_group(sub_shape, translation_dict)

        except Exception as e:
            logging.info(f" 处理组合形状失败: {e}")

        return replaced_count

    def _final_font_size_fix(self, prs, translation_dict: Dict[str, str]):
        """最后的字体大小修复，确保所有翻译文本都应用了字体调整"""
        try:
            fixed_count = 0
            for slide_idx, slide in enumerate(prs.slides):
                logging.info(f"🔍 检查幻灯片 {slide_idx + 1}...")

                # 检查占位符
                try:
                    for placeholder in slide.placeholders:
                        if hasattr(placeholder, 'text_frame') and placeholder.text_frame:
                            for paragraph in placeholder.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if run.text.strip() and run.text.strip() in [v for v in translation_dict.values() if
                                                                                 v]:
                                        # 这是翻译后的文本，检查字体大小
                                        if self._check_and_fix_font_size(run):
                                            fixed_count += 1
                except Exception as e:
                    logging.info(f" 检查占位符失败: {e}")

                # 检查所有形状
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, 'text_frame') and shape.text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if run.text.strip() and run.text.strip() in [v for v in translation_dict.values() if
                                                                                 v]:
                                        if self._check_and_fix_font_size(run):
                                            fixed_count += 1

                        # 检查表格
                        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                            table = shape.table
                            for row in table.rows:
                                for cell in row.cells:
                                    if hasattr(cell, 'text_frame') and cell.text_frame:
                                        for paragraph in cell.text_frame.paragraphs:
                                            for run in paragraph.runs:
                                                if run.text.strip() and run.text.strip() in [v for v in
                                                                                             translation_dict.values()
                                                                                             if v]:
                                                    if self._check_and_fix_font_size(run):
                                                        fixed_count += 1

                        # 检查组合形状
                        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                            fixed_count += self._fix_group_font_sizes(shape, translation_dict)

                    except Exception as e:
                        logging.info(f" 检查形状失败: {e}")
                        continue

            if fixed_count > 0:
                logging.info(f" 最后修复了 {fixed_count} 个字体大小")

        except Exception as e:
            logging.info(f" 最后字体检查失败: {e}")

    def _check_and_fix_font_size(self, run) -> bool:
        """检查并修复单个run的字体大小"""
        try:
            if run.font.size is None:
                # 没有设置字体大小，应用默认调整
                from pptx.util import Pt
                default_size = 18  # 默认字体大小
                adjusted_size = max(8, min(int(default_size * self.font_size_adjustment), default_size))
                run.font.size = Pt(adjusted_size)
                logging.info(f"修复缺失字体大小: 设置为 {adjusted_size}pt")
                return True

            # 检查字体大小是否过大（可能未被调整）
            current_size = run.font.size.pt if hasattr(run.font.size, 'pt') else run.font.size
            expected_max_size = 50 * self.font_size_adjustment  # 假设最大原始字体为50pt

            if current_size > expected_max_size:
                # 字体可能未被调整，重新调整
                adjusted_size = max(8, min(int(current_size * self.font_size_adjustment), int(current_size)))
                from pptx.util import Pt
                run.font.size = Pt(adjusted_size)
                logging.info(f"重新调整字体大小: {current_size}pt → {adjusted_size}pt")
                return True

            return False

        except Exception as e:
            logging.info(f" 检查字体大小失败: {e}")
            return False

    def _fix_group_font_sizes(self, group_shape, translation_dict: Dict[str, str]) -> int:
        """修复组合形状中的字体大小"""
        fixed_count = 0

        try:
            for sub_shape in group_shape.shapes:
                if hasattr(sub_shape, 'text_frame') and sub_shape.text_frame:
                    for paragraph in sub_shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text.strip() and run.text.strip() in [v for v in translation_dict.values() if v]:
                                if self._check_and_fix_font_size(run):
                                    fixed_count += 1

                # 递归处理嵌套组合
                elif sub_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    fixed_count += self._fix_group_font_sizes(sub_shape, translation_dict)

        except Exception as e:
            logging.info(f" 修复组合字体失败: {e}")

        return fixed_count

    def _replace_text_in_group_enhanced(self, group_shape, translation_dict: Dict[str, str]) -> int:
        """增强版组合形状文本替换，处理每个run"""
        replaced_count = 0

        try:
            for sub_shape in group_shape.shapes:
                if hasattr(sub_shape, 'text_frame') and sub_shape.text_frame:
                    for paragraph in sub_shape.text_frame.paragraphs:
                        # 先处理段落中的每个run
                        for run in paragraph.runs:
                            if run.text.strip():
                                original_text = run.text.strip()
                                if original_text in translation_dict:
                                    translated_text = translation_dict[original_text]
                                    if original_text != translated_text:
                                        try:
                                            format_info = _extract_format_info(run)
                                            run.text = translated_text
                                            self._apply_format(run, format_info)
                                            replaced_count += 1
                                            logging.info(f"备用组合run替换: '{original_text}' → '{translated_text}'")
                                        except Exception as e:
                                            logging.info(f"备用组合run替换失败: {original_text} - {e}")

                        # 再处理整个段落
                        if paragraph.text.strip():
                            original_text = paragraph.text.strip()
                            if original_text in translation_dict:
                                translated_text = translation_dict[original_text]
                                if original_text != translated_text:
                                    try:
                                        if paragraph.runs:
                                            first_run = paragraph.runs[0]
                                            format_info = _extract_format_info(first_run)

                                            paragraph.clear()
                                            new_run = paragraph.add_run()
                                            new_run.text = translated_text
                                            self._apply_format(new_run, format_info)
                                            replaced_count += 1
                                            logging.info(f"🔄 备用组合段落替换: '{original_text}' → '{translated_text}'")
                                    except Exception as e:
                                        logging.warning(f"备用组合段落替换失败: {original_text} - {e}")

                # 递归处理嵌套组合
                elif sub_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    replaced_count += self._replace_text_in_group_enhanced(sub_shape, translation_dict)

        except Exception as e:
            logging.warning(f"处理增强组合形状失败: {e}")

        return replaced_count

    def _fallback_smartart_replace(self, prs, item, translated_text, original_text):
        """SmartArt文本替换的备用方法"""
        try:
            slide = prs.slides[item['slide_index']]
            shape = slide.shapes[item['shape_index']]

            # 尝试在形状的所有文本框中查找并替换
            if hasattr(shape, 'text_frame') and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if original_text in paragraph.text:
                        # 找到包含目标文本的段落
                        if paragraph.runs:
                            # 保存格式并替换整个段落
                            main_format = _extract_format_info(paragraph.runs[0])
                            new_text = paragraph.text.replace(original_text, translated_text)
                            paragraph.clear()
                            run = paragraph.add_run()
                            run.text = new_text
                            self._apply_format(run, main_format)
                            logging.info(f"🔄 SmartArt备用替换成功: {original_text} → {translated_text}")
                            return
                        else:
                            # 直接替换段落文本
                            paragraph.text = paragraph.text.replace(original_text, translated_text)
                            logging.info(f"🔄 SmartArt备用替换成功: {original_text} → {translated_text}")
                            return

            logging.warning(f"SmartArt备用替换失败: {original_text}")

        except Exception as e:
            logging.warning(f"SmartArt备用替换错误: {original_text} - {e}")

    def _replace_smartart_content_in_pptx(self, pptx_path: str, output_path: str, translation_dict: Dict[str, str]):
        """直接修改PPTX内部的SmartArt XML文件"""
        import tempfile
        import shutil

        try:
            logging.info("🔧 直接修改SmartArt XML内容...")

            # 创建临时目录
            with tempfile.TemporaryDirectory() as temp_dir:
                # 复制原文件到临时位置
                temp_pptx = os.path.join(temp_dir, "temp.pptx")
                shutil.copy2(pptx_path, temp_pptx)

                # 解压PPTX文件
                extract_dir = os.path.join(temp_dir, "extracted")
                with zipfile.ZipFile(temp_pptx, 'r') as zip_file:
                    zip_file.extractall(extract_dir)

                # 修改diagram数据文件
                diagrams_dir = os.path.join(extract_dir, "ppt", "diagrams")
                if os.path.exists(diagrams_dir):
                    data_files = [f for f in os.listdir(diagrams_dir) if f.startswith('data') and f.endswith('.xml')]

                    replaced_count = 0
                    for data_file in data_files:
                        data_path = os.path.join(diagrams_dir, data_file)
                        if self._modify_diagram_xml(data_path, translation_dict):
                            replaced_count += 1

                    logging.info(f"修改了 {replaced_count} 个diagram数据文件")

                # 重新打包PPTX文件
                _repack_pptx(extract_dir, output_path)
                logging.info(f" SmartArt内容已写入: {output_path}")

        except Exception as e:
            logging.warning(f"SmartArt内容写入失败: {e}")
            # 如果失败，至少复制原文件
            try:
                shutil.copy2(pptx_path, output_path)
            except:
                pass

    def _modify_diagram_xml(self, xml_path: str, translation_dict: Dict[str, str]) -> bool:
        """修改单个diagram XML文件，保持字体格式并调整所有字体大小"""
        try:
            import xml.etree.ElementTree as ET
            # 解析XML文件
            tree = ET.parse(xml_path)
            root = tree.getroot()

            # 命名空间映射
            namespaces = {
                'dgm': 'http://schemas.openxmlformats.org/drawingml/2006/diagram',
                'a':   'http://schemas.openxmlformats.org/drawingml/2006/main'
            }

            # 1. 批量调整所有字体大小
            self._adjust_all_font_sizes(root)

            # 2. 替换 diagram 文本节点
            for t in root.findall('.//dgm:t', namespaces):
                orig = t.text.strip() if t.text else None
                if orig and orig in translation_dict:
                    new_text = translation_dict[orig]
                    # 检查是否是中文翻英文
                    is_zh_to_en = (
                        any('\u4e00' <= ch <= '\u9fff' for ch in orig) and
                        any(ch.isalpha() and ord(ch) < 256 for ch in new_text)
                    )
                    if is_zh_to_en:
                        # 调整字体以适应英文
                        _adjust_font_for_translation(t, new_text, root)
                    else:
                        t.text = new_text

            # 写回 XML，保留声明
            tree.write(xml_path, encoding='utf-8', xml_declaration=True)
            return True

        except Exception as e:
            logging.warning(f"修改 diagram XML ({xml_path}) 失败: {e}")
            return False


    def _adjust_all_font_sizes(self, root):
        """调整XML中所有的字体大小"""
        try:
            adjusted_count = 0
            for elem in root.iter():
                if elem.tag.endswith('}rPr'):
                    current_size = elem.get('sz')
                    if current_size:
                        # 将字符串转换为整数，应用调整比例（只缩不放）
                        original_size = int(current_size)
                        new_size = min(int(original_size * self.font_size_adjustment), original_size)
                        elem.set('sz', str(new_size))
                        adjusted_count += 1

            if adjusted_count > 0:
                logging.info(f"调整了 {adjusted_count} 个字体大小 (比例: {self.font_size_adjustment})")
        except Exception as e:
            logging.warning(f"批量字体大小调整失败: {e}")

    def _adjust_font_size(self, rPr_elem):
        """调整字体大小"""
        try:
            current_size = rPr_elem.get('sz')
            if current_size:
                # 将字符串转换为整数，应用调整比例（只缩不放）
                original_size = int(current_size)
                new_size = min(int(original_size * self.font_size_adjustment), original_size)
                rPr_elem.set('sz', str(new_size))
                logging.info(f"字体大小调整: {original_size} → {new_size} (比例: {self.font_size_adjustment})")
        except Exception as e:
            logging.warning(f"字体大小调整失败: {e}")

    def _adjust_xml_font_size(self, xml_element):
        """调整XML元素中的字体大小"""
        try:
            # 查找父元素或兄弟元素中的字体大小设置
            parent = xml_element.getparent() if hasattr(xml_element, 'getparent') else None

            if parent is not None:
                # 查找rPr元素（包含字体属性）
                for elem in parent.iter():
                    if elem.tag.endswith('}rPr') or 'rPr' in elem.tag:
                        sz_attr = elem.get('sz')
                        if sz_attr:
                            try:
                                original_size = int(sz_attr)
                                adjusted_size = min(int(original_size * self.font_size_adjustment), original_size)
                                elem.set('sz', str(adjusted_size))
                                logging.info(
                                    f"XML字体大小调整: {original_size} → {adjusted_size} (比例: {self.font_size_adjustment})")
                            except (ValueError, TypeError):
                                pass

                        # 查找字体大小相关的其他属性
                        for child in elem:
                            if hasattr(child, 'tag') and ('sz' in child.tag or 'size' in child.tag.lower()):
                                size_val = child.get('val') or child.text
                                if size_val:
                                    try:
                                        original_size = int(size_val)
                                        adjusted_size = min(int(original_size * self.font_size_adjustment), original_size)
                                        if child.get('val'):
                                            child.set('val', str(adjusted_size))
                                        else:
                                            child.text = str(adjusted_size)
                                        logging.info(f"XML子元素字体调整: {original_size} → {adjusted_size}")
                                    except (ValueError, TypeError):
                                        pass

        except Exception as e:
            logging.warning(f"XML字体大小调整失败: {e}")

    def replace_text_in_ppt_optimized(self, ppt_path: str, output_path: str,
                                      text_data: List[Dict], translation_dict: Dict[str, str]):
        """优化的PPT文本替换函数"""
        prs = Presentation(ppt_path)
        replaced_count = 0
        
        # 按优先级排序：run级别 > 段落级别 > 整体级别
        def get_priority(item):
            if item.get('is_table_run'):
                return 1  # 最高优先级
            elif item.get('is_table_paragraph'):
                return 2
            elif item.get('is_cell_full_text'):
                return 5  # 较低优先级
            elif item.get('is_run'):
                return 1
            elif 'paragraph_obj' in item:
                return 3
            else:
                return 4
        
        # 排序文本数据以确保精确替换优先
        sorted_text_data = sorted(text_data, key=get_priority)
        
        # 跟踪已处理的对象，避免重复替换
        processed_runs = set()
        processed_paragraphs = set()
        processed_cells = set()
        
        for item in sorted_text_data:
            original_text = item['text']
            
            # 尝试精确匹配，然后尝试标准化匹配
            translated_text = translation_dict.get(original_text)
            if translated_text is None:
                normalized_text = normalize_text(original_text)
                translated_text = translation_dict.get(normalized_text, original_text)
            
            if original_text == translated_text:
                continue
                
            try:
                # 处理长文本的智能换行
                if len(translated_text) > 50:
                    translated_text = self._smart_text_wrapping(translated_text)
                
                # 处理占位符文本
                if item.get('is_placeholder'):
                    slide = prs.slides[item['slide_index']]
                    placeholder = slide.placeholders[item['placeholder_idx']]
                    
                    # 优化文本框以支持自动换行
                    self._optimize_textframe_for_wrapping(
                        placeholder.text_frame, 
                        original_text, 
                        translated_text
                    )
                    
                    # 处理段落
                    paragraph = placeholder.text_frame.paragraphs[item['paragraph_index']]
                    
                    if paragraph.runs:
                        # 保存格式并应用智能大小调整
                        format_info = _extract_format_info(paragraph.runs[0])
                        paragraph.clear()
                        run = paragraph.add_run()
                        run.text = translated_text
                        self._apply_format_with_smart_sizing(
                            run, format_info, original_text, translated_text
                        )
                        # 应用智能对齐优化
                        self._apply_smart_alignment_to_paragraph(paragraph, original_text, translated_text)
                        # 设置行间距
                        self._set_paragraph_line_spacing(paragraph, 0.9)
                    
                    replaced_count += 1
                    logging.info(f"占位符替换: '{original_text}' → '{translated_text}'")
                
                # 处理表格单元格的run级别替换（最优先，最精确）
                elif item.get('is_table_run') and 'run_obj' in item:
                    try:
                        run = item['run_obj']
                        run_id = id(run)
                        
                        # 检查是否已处理过这个run对象
                        if run_id in processed_runs:
                            continue
                        processed_runs.add(run_id)
                        
                        # 保存原始格式信息
                        format_info = _extract_format_info(run)
                        # 替换文本
                        run.text = translated_text
                        # 应用格式（包括调整后的字体大小）
                        self._apply_format_with_smart_sizing(
                            run, format_info, original_text, translated_text
                        )
                        replaced_count += 1
                        logging.info(f"表格run替换: '{original_text}' → '{translated_text}'")
                    except Exception as e:
                        logging.warning(f"表格run替换失败: {original_text} - {e}")
                        continue
                
                # 处理表格单元格的段落级别替换
                elif item.get('is_table_paragraph') and 'paragraph_obj' in item:
                    try:
                        paragraph = item['paragraph_obj']
                        para_id = id(paragraph)
                        
                        # 检查是否已处理过这个段落对象
                        if para_id in processed_paragraphs:
                            continue
                        processed_paragraphs.add(para_id)
                        
                        if paragraph.runs:
                            format_info = _extract_format_info(paragraph.runs[0])
                            paragraph.clear()
                            run = paragraph.add_run()
                            run.text = translated_text
                            self._apply_format_with_smart_sizing(
                                run, format_info, original_text, translated_text
                            )
                        else:
                            run = paragraph.add_run()
                            run.text = translated_text
                            from pptx.util import Pt
                            run.font.size = Pt(max(8, int(12 * self.font_size_adjustment)))
                        
                        # 应用智能对齐优化
                        self._apply_smart_alignment_to_paragraph(paragraph, original_text, translated_text)
                        # 设置行间距（表格内文本稍微紧密）
                        self._set_paragraph_line_spacing(paragraph, 0.8)
                        
                        replaced_count += 1
                        logging.info(f"表格段落替换: '{original_text}' → '{translated_text}'")
                    except Exception as e:
                        logging.warning(f"表格段落替换失败: {original_text} - {e}")
                        continue
                
                # 处理表格单元格整体替换（备用）
                elif item.get('is_cell_full_text') or ('cell_obj' in item or ('cell_row' in item and 'cell_col' in item)):
                    try:
                        slide = prs.slides[item['slide_index']]
                        shape = slide.shapes[item['shape_index']]
                        
                        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                            table = shape.table
                            cell = table.rows[item['cell_row']].cells[item['cell_col']]
                            cell_id = f"{item['slide_index']}_{item['shape_index']}_{item['cell_row']}_{item['cell_col']}"
                            
                            # 检查是否已处理过这个单元格
                            if cell_id in processed_cells:
                                continue
                            processed_cells.add(cell_id)
                            
                            # 处理表格单元格中的文本
                            if hasattr(cell, 'text_frame') and cell.text_frame:
                                # 如果单元格有text_frame，处理段落和runs
                                if cell.text_frame.paragraphs:
                                    paragraph = cell.text_frame.paragraphs[0]  # 使用第一个段落
                                    if paragraph.runs:
                                        format_info = _extract_format_info(paragraph.runs[0])
                                        paragraph.clear()
                                        run = paragraph.add_run()
                                        run.text = translated_text
                                        self._apply_format_with_smart_sizing(
                                            run, format_info, original_text, translated_text
                                        )
                                    else:
                                        # 没有runs，直接添加
                                        run = paragraph.add_run()
                                        run.text = translated_text
                                        # 应用默认表格字体调整
                                        from pptx.util import Pt
                                        run.font.size = Pt(max(8, int(12 * self.font_size_adjustment)))
                                else:
                                    # 没有段落，创建一个
                                    paragraph = cell.text_frame.add_paragraph()
                                    run = paragraph.add_run()
                                    run.text = translated_text
                                    from pptx.util import Pt
                                    run.font.size = Pt(max(8, int(12 * self.font_size_adjustment)))
                            else:
                                # 直接设置单元格文本（备用方法）
                                cell.text = translated_text
                            
                            replaced_count += 1
                            logging.info(f"表格整体替换: '{original_text}' → '{translated_text}'")
                        else:
                            logging.warning(f"预期为表格但类型不匹配: {shape.shape_type}")
                            continue
                            
                    except (IndexError, AttributeError) as e:
                        logging.warning(f"表格文本替换失败: {original_text} - {e}")
                        continue
                
                # 处理普通文本框
                elif 'paragraph_obj' in item:
                    slide = prs.slides[item['slide_index']]
                    shape = slide.shapes[item['shape_index']]
                    
                    # 优化文本框以支持自动换行
                    if hasattr(shape, 'text_frame'):
                        self._optimize_textframe_for_wrapping(
                            shape.text_frame, 
                            original_text, 
                            translated_text
                        )
                    
                    paragraph = shape.text_frame.paragraphs[item['paragraph_index']]
                    
                    if paragraph.runs:
                        format_info = _extract_format_info(paragraph.runs[0])
                        paragraph.clear()
                        run = paragraph.add_run()
                        run.text = translated_text
                        self._apply_format_with_smart_sizing(
                            run, format_info, original_text, translated_text
                        )
                        # 应用智能对齐优化
                        self._apply_smart_alignment_to_paragraph(paragraph, original_text, translated_text)
                        # 设置行间距
                        self._set_paragraph_line_spacing(paragraph, 0.9)
                    
                    replaced_count += 1
                    logging.info(f"文本框替换: '{original_text}' → '{translated_text}'")
                
            except Exception as e:
                logging.warning(f"替换失败: {original_text} - {str(e)}")
                continue
        
        # 如果替换数量太少，使用备用的简单替换策略
        expected_replacements = len([item for item in text_data if translation_dict.get(item['text'], item['text']) != item['text']])
        if expected_replacements > 0 and replaced_count < expected_replacements * 0.5:
            logging.info(f"替换数量偏少 ({replaced_count}/{expected_replacements})，使用备用替换策略...")
            
            # 调试：列出未被替换的文本
            untranslated_texts = []
            for item in text_data:
                text = item['text']
                if text in translation_dict and translation_dict[text] != text:
                    # 这是应该被翻译的文本，检查是否包含中文
                    if any('\u4e00' <= char <= '\u9fff' for char in text):
                        untranslated_texts.append(text)
            
            if untranslated_texts:
                logging.warning(f"可能未完全翻译的中文文本样例: {untranslated_texts[:5]}")
            
            backup_count = self._simple_text_replacement(prs, translation_dict)
            replaced_count += backup_count
            logging.info(f"备用策略增加 {backup_count} 处替换")
            
            # 如果仍然有未翻译的文本，使用强力替换
            if backup_count == 0 and untranslated_texts:
                logging.info("执行强力文本替换...")
                force_count = self._force_text_replacement(prs, translation_dict)
                replaced_count += force_count
                logging.info(f"强力替换增加 {force_count} 处替换")
        
        prs.save(output_path)
        logging.info(f"PPT翻译完成！共替换 {replaced_count} 处文本")

    def _post_process_optimization(self, pptx_path: str):
        """后处理优化，确保所有文本显示正确并启用自动换行"""
        from pptx.util import Pt
        from pptx.enum.text import MSO_AUTO_SIZE
        
        prs = Presentation(pptx_path)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    # 确保所有文本框都启用了自动换行
                    try:
                        shape.text_frame.word_wrap = True
                        logging.debug(f"为形状启用自动换行")
                    except Exception as e:
                        logging.debug(f"启用自动换行失败: {e}")
                    
                    # 检查文本是否溢出
                    if self._is_text_overflow(shape):
                        # 进一步缩小字体
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.font.size:
                                    current_size = run.font.size.pt
                                    new_size = max(8, int(current_size * 0.9))
                                    run.font.size = Pt(new_size)
                        
                        # 再次设置自动调整
                        shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                
                # 处理表格中的文本框
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            if hasattr(cell, 'text_frame') and cell.text_frame:
                                try:
                                    cell.text_frame.word_wrap = True
                                    logging.debug(f"为表格单元格启用自动换行")
                                except Exception as e:
                                    logging.debug(f"表格单元格换行设置失败: {e}")
        
        prs.save(pptx_path)
        logging.info("后处理优化完成")

    def translate_ppt(self, ppt_path: str, output_path: str):
        """
        翻译PPT文件
        """
        logging.info(f"开始翻译PPT文件: {ppt_path}")

        try:
            # 1. 提取文本
            logging.info("步骤1: 提取文本...")
            text_data = self.extract_text_from_ppt(ppt_path)

            if not text_data:
                logging.warning("未找到需要翻译的文本")
                return

            logging.info(f"共找到 {len(text_data)} 段文本")

            # 2. 去重并翻译（使用标准化文本）
            logging.info("步骤2: 去重文本...")
            
            # 创建标准化文本到原文本的映射
            normalized_to_original = {}
            for item in text_data:
                original = item['text']
                normalized = normalize_text(original)
                if should_translate(normalized) and normalized not in normalized_to_original:
                    normalized_to_original[normalized] = original
            
            unique_texts = list(normalized_to_original.values())
            logging.info(f"去重后需翻译 {len(unique_texts)} 段文本")

            logging.info("步骤3: 批量翻译...")
            
            # 调试：显示去重后的文本样例
            if len(unique_texts) > 0:
                sample_texts = unique_texts[:3]
                logging.info(f"待翻译文本样例: {sample_texts}")
            
            translated_texts = self.translator.translate_text_batch(unique_texts)
            translation_dict = dict(zip(unique_texts, translated_texts))
            
            # 创建标准化版本的翻译字典，用于更好的匹配
            normalized_translation_dict = {}
            for original, translated in translation_dict.items():
                normalized_original = normalize_text(original)
                normalized_translation_dict[normalized_original] = translated
                # 同时保持原版本
                normalized_translation_dict[original] = translated

            # 3. 使用优化的智能替换文本（传递标准化字典）
            logging.info("步骤4: 智能替换文本...")
            self.replace_text_in_ppt_optimized(ppt_path, output_path, text_data, normalized_translation_dict)

            # 4. 直接修改SmartArt内容
            logging.info("步骤5: 修改SmartArt内容...")
            self._replace_smartart_content_in_pptx(output_path, output_path, translation_dict)

            # 5. 后处理优化
            logging.info("步骤6: 后处理优化...")
            self._post_process_optimization(output_path)

            # 6. 生成SmartArt翻译参考文件
            logging.info("步骤7: 生成SmartArt翻译参考...")
            _generate_smartart_reference(text_data, translation_dict, output_path)

            # 7. 最终字体大小限制（使用更智能的限制）
            _force_cap_font(output_path, max_title_pt=28, max_body_pt=20)

            logging.info(" 翻译完成！")

        except KeyboardInterrupt:
            logging.error("\n 用户中断翻译过程")
        except Exception as e:
            logging.error(f" 翻译过程出错: {e}")
            import traceback
            logging.error(traceback.format_exc())
    
    def _force_text_replacement(self, prs, translation_dict: Dict[str, str]) -> int:
        """强力文本替换，使用字符串匹配和替换"""
        replaced_count = 0
        
        for slide in prs.slides:
            for shape in slide.shapes:
                try:
                    # 处理所有文本框类型
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        replaced_count += self._force_replace_in_textframe(shape.text_frame, translation_dict)
                    
                    # 处理表格
                    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        table = shape.table
                        for row in table.rows:
                            for cell in row.cells:
                                if hasattr(cell, 'text_frame') and cell.text_frame:
                                    replaced_count += self._force_replace_in_textframe(cell.text_frame, translation_dict)
                                # 直接处理单元格文本
                                elif hasattr(cell, 'text') and cell.text:
                                    for original, translated in translation_dict.items():
                                        if original != translated and original in cell.text:
                                            cell.text = cell.text.replace(original, translated)
                                            replaced_count += 1
                                            logging.info(f"强力单元格替换: '{original}' → '{translated}'")
                    
                    # 处理组合形状
                    elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        replaced_count += self._force_replace_in_group(shape, translation_dict)
                        
                except Exception as e:
                    logging.debug(f"强力替换处理形状时出错: {e}")
                    continue
        
        return replaced_count
    
    def _force_replace_in_textframe(self, text_frame, translation_dict: Dict[str, str]) -> int:
        """在文本框中进行强力替换"""
        replaced_count = 0
        
        try:
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text:
                        original_text = run.text
                        new_text = original_text
                        
                        # 对每个翻译对进行字符串替换
                        for original, translated in translation_dict.items():
                            if original != translated and original in new_text:
                                new_text = new_text.replace(original, translated)
                                replaced_count += 1
                                logging.info(f"强力run替换: '{original}' → '{translated}'")
                        
                        if new_text != original_text:
                            run.text = new_text
                            
        except Exception as e:
            logging.debug(f"强力文本框替换出错: {e}")
        
        return replaced_count
    
    def _force_replace_in_group(self, group_shape, translation_dict: Dict[str, str]) -> int:
        """在组合形状中进行强力替换"""
        replaced_count = 0
        
        try:
            if hasattr(group_shape, 'shapes'):
                for sub_shape in group_shape.shapes:
                    if hasattr(sub_shape, 'text_frame') and sub_shape.text_frame:
                        replaced_count += self._force_replace_in_textframe(sub_shape.text_frame, translation_dict)
                    elif sub_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        replaced_count += self._force_replace_in_group(sub_shape, translation_dict)
        except Exception as e:
            logging.debug(f"强力组合形状替换出错: {e}")
        
        return replaced_count


def main():
    translator = PPTTranslator()
    translator.translate_ppt(
        ppt_path="0725.pptx",
        output_path="0725_translated_en.pptx"# 翻译成英文
    )

if __name__ == "__main__":
    main()