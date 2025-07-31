"""
PyTest configuration and fixtures for Offitrans tests
"""

import pytest
import tempfile
import shutil
from pathlib import Path
from typing import Generator

from offitrans.core.config import Config
from offitrans.translators.google import GoogleTranslator
from offitrans.processors.excel import ExcelProcessor


@pytest.fixture
def temp_dir() -> Generator[Path, None, None]:
    """
    Create a temporary directory for tests.
    
    Yields:
        Path to temporary directory
    """
    temp_path = Path(tempfile.mkdtemp())
    try:
        yield temp_path
    finally:
        shutil.rmtree(temp_path, ignore_errors=True)


@pytest.fixture
def sample_text_data():
    """
    Sample text data for testing translations.
    
    Returns:
        List of sample texts in different languages
    """
    return [
        "Hello, world!",
        "Hello world!",  # Changed from Chinese to English for consistency
        "สวัสดีชาวโลก!",
        "Bonjour le monde!",
        "Hola mundo!",
        "123",  # Should not be translated
        "test@example.com",  # Should not be translated
        "",  # Empty string
        "   ",  # Whitespace only
    ]


@pytest.fixture
def mock_translator():
    """
    Create a mock translator for testing.
    
    Returns:
        Mock translator instance
    """
    class MockTranslator:
        def __init__(self, source_lang="auto", target_lang="en"):
            self.source_lang = source_lang
            self.target_lang = target_lang
        
        def translate_text(self, text: str) -> str:
            # Simple mock translation - just add prefix
            if text.strip():
                return f"[TRANSLATED_{self.target_lang}] {text}"
            return text
        
        def translate_text_batch(self, texts):
            return [self.translate_text(text) for text in texts]
    
    return MockTranslator()


@pytest.fixture
def test_config():
    """
    Create a test configuration.
    
    Returns:
        Config instance for testing
    """
    config = Config()
    # Use conservative settings for tests
    config.translator.max_workers = 2
    config.translator.timeout = 10
    config.translator.retry_count = 1
    config.cache.enabled = False  # Disable cache for consistent tests
    return config


@pytest.fixture
def google_translator():
    """
    Create a Google translator instance for testing.
    Note: This will use the free API without API key for basic tests.
    
    Returns:
        GoogleTranslator instance
    """
    return GoogleTranslator(
        source_lang="auto",
        target_lang="en",
        use_free_api=True,
        max_workers=1,  # Conservative for tests
        timeout=10
    )


@pytest.fixture
def excel_processor(test_config):
    """
    Create an Excel processor for testing.
    
    Args:
        test_config: Test configuration fixture
        
    Returns:
        ExcelProcessor instance
    """
    try:
        return ExcelProcessor(config=test_config)
    except ImportError:
        pytest.skip("openpyxl not available")


class TestFileGenerator:
    """Helper class to generate test files"""
    
    @staticmethod
    def create_simple_excel(file_path: Path, texts: list = None):
        """Create a simple Excel file for testing"""
        try:
            from openpyxl import Workbook
            
            wb = Workbook()
            ws = wb.active
            ws.title = "Test Sheet"
            
            test_texts = texts or [
                "Hello World",
                "Hello World",
                "Test Data",
                "123",
                "sample@email.com"
            ]
            
            for i, text in enumerate(test_texts, 1):
                ws[f'A{i}'] = text
            
            wb.save(file_path)
            return True
            
        except ImportError:
            return False
    
    @staticmethod
    def create_simple_word(file_path: Path, texts: list = None):
        """Create a simple Word document for testing"""
        try:
            from docx import Document
            
            doc = Document()
            
            test_texts = texts or [
                "This is a test document",
                "This is a test document",
                "Sample paragraph with formatting"
            ]
            
            for text in test_texts:
                doc.add_paragraph(text)
            
            doc.save(file_path)
            return True
            
        except ImportError:
            return False
    
    @staticmethod
    def create_simple_ppt(file_path: Path, texts: list = None):
        """Create a simple PowerPoint for testing"""
        try:
            from pptx import Presentation
            
            prs = Presentation()
            
            test_texts = texts or [
                "Test Slide Title",
                "Test Slide Content",
                "Sample content"
            ]
            
            for text in test_texts:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = text
            
            prs.save(file_path)
            return True
            
        except ImportError:
            return False


@pytest.fixture
def file_generator():
    """
    File generator fixture for creating test files.
    
    Returns:
        TestFileGenerator instance
    """
    return TestFileGenerator()


# Pytest configuration
def pytest_configure(config):
    """Configure pytest with custom markers"""
    config.addinivalue_line("markers", "integration: mark test as integration test")
    config.addinivalue_line("markers", "slow: mark test as slow running") 
    config.addinivalue_line("markers", "requires_api: mark test as requiring API access")
    config.addinivalue_line("markers", "requires_openpyxl: mark test as requiring openpyxl")
    config.addinivalue_line("markers", "requires_docx: mark test as requiring python-docx")
    config.addinivalue_line("markers", "requires_pptx: mark test as requiring python-pptx")


def pytest_collection_modifyitems(config, items):
    """Modify test collection to add markers based on test names and paths"""
    for item in items:
        # Mark integration tests
        if "integration" in str(item.fspath):
            item.add_marker(pytest.mark.integration)
        
        # Mark tests that require specific libraries
        if "excel" in str(item.fspath).lower():
            item.add_marker(pytest.mark.requires_openpyxl)
        elif "word" in str(item.fspath).lower():
            item.add_marker(pytest.mark.requires_docx)
        elif "powerpoint" in str(item.fspath).lower() or "ppt" in str(item.fspath).lower():
            item.add_marker(pytest.mark.requires_pptx)
        
        # Mark API tests
        if "api" in item.name.lower() or "translate" in item.name.lower():
            item.add_marker(pytest.mark.requires_api)