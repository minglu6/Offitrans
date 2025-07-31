"""
PowerPoint file processor for Offitrans

This module provides functionality to translate PowerPoint presentations
while preserving layout and formatting.
"""

import logging
from typing import List, Dict, Any
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False

from .base import BaseProcessor
from ..exceptions.errors import PowerPointProcessorError

logger = logging.getLogger(__name__)


class PowerPointProcessor(BaseProcessor):
    """
    PowerPoint file processor that handles translation while preserving layout.
    
    This processor can handle:
    - Slide text content
    - Text boxes and shapes
    - Title and content placeholders
    - Text formatting
    """
    
    def __init__(self, **kwargs):
        """
        Initialize PowerPoint processor.
        
        Args:
            **kwargs: Additional arguments passed to BaseProcessor
        """
        if not PYTHON_PPTX_AVAILABLE:
            raise PowerPointProcessorError(
                "python-pptx library is required for PowerPoint processing",
                details="Install with: pip install python-pptx"
            )
        
        super().__init__(**kwargs)
    
    def supports_file_type(self, file_path: str) -> bool:
        """
        Check if file type is supported.
        
        Args:
            file_path: Path to the file
            
        Returns:
            True if file type is supported
        """
        supported_extensions = {'.pptx', '.ppt'}
        return Path(file_path).suffix.lower() in supported_extensions
    
    def extract_text(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Extract text content from PowerPoint presentation.
        
        Args:
            file_path: Path to the PowerPoint file
            
        Returns:
            List of dictionaries containing text and metadata
        """
        text_data = []
        
        try:
            prs = Presentation(file_path)
            logger.info(f"Successfully opened PowerPoint file: {file_path}")
            logger.info(f"Presentation has {len(prs.slides)} slides")
            
            for slide_idx, slide in enumerate(prs.slides):
                logger.debug(f"Processing slide {slide_idx + 1}")
                
                # Extract text from shapes
                for shape_idx, shape in enumerate(slide.shapes):
                    if hasattr(shape, "text") and shape.text.strip():
                        # Get shape type and properties
                        shape_info = self._extract_shape_info(shape)
                        
                        text_data.append({
                            'text': shape.text,
                            'slide_index': slide_idx,
                            'shape_index': shape_idx,
                            'shape_info': shape_info,
                            'type': 'shape_text'
                        })
                        
                        logger.debug(f"Extracted text from slide {slide_idx + 1}, shape {shape_idx}: '{shape.text[:50]}...'")
                    
                    # Extract text from text frames within shapes
                    if hasattr(shape, "text_frame"):
                        for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                            if paragraph.text.strip():
                                para_info = self._extract_paragraph_info(paragraph)
                                
                                text_data.append({
                                    'text': paragraph.text,
                                    'slide_index': slide_idx,
                                    'shape_index': shape_idx,
                                    'paragraph_index': para_idx,
                                    'paragraph_info': para_info,
                                    'type': 'paragraph_text'
                                })
                                
                                logger.debug(f"Extracted paragraph from slide {slide_idx + 1}, shape {shape_idx}, para {para_idx}: '{paragraph.text[:50]}...'")
            
            logger.info(f"Total extracted {len(text_data)} text elements from PowerPoint")
            return text_data
            
        except Exception as e:
            raise PowerPointProcessorError(
                f"Failed to extract text from PowerPoint file",
                details=str(e),
                file_path=file_path
            ) from e
    
    def translate_and_save(self, 
                          file_path: str, 
                          output_path: str,
                          target_language: str = "en") -> bool:
        """
        Translate PowerPoint presentation and save to output path.
        
        Args:
            file_path: Path to input PowerPoint file
            output_path: Path for output PowerPoint file
            target_language: Target language code
            
        Returns:
            True if successful, False otherwise
        """
        try:
            # Step 1: Extract text and metadata
            logger.info("Step 1: Extracting text from PowerPoint presentation...")
            text_data = self.extract_text(file_path)
            
            if not text_data:
                logger.warning("No translatable text found in PowerPoint presentation")
                return False
            
            # Step 2: Preprocess and translate texts
            logger.info("Step 2: Translating texts...")
            original_texts = [item['text'] for item in text_data]
            unique_texts, metadata = self.preprocess_texts(original_texts)
            translated_unique = self.translate_texts(unique_texts, target_language)
            translated_texts = self.postprocess_translations(original_texts, translated_unique, metadata)
            
            # Step 3: Apply translations to PowerPoint presentation
            logger.info("Step 3: Applying translations to PowerPoint presentation...")
            success = self._replace_text_with_format(
                file_path, output_path, text_data, translated_texts, target_language
            )
            
            if success:
                logger.info(f"Successfully translated PowerPoint presentation: {output_path}")
                return True
            else:
                logger.error("Failed to apply translations to PowerPoint presentation")
                return False
                
        except Exception as e:
            logger.error(f"Error translating PowerPoint presentation: {e}")
            return False
    
    def _replace_text_with_format(self, 
                                 ppt_path: str, 
                                 output_path: str,
                                 text_data: List[Dict[str, Any]], 
                                 translated_texts: List[str],
                                 target_language: str = "en") -> bool:
        """
        Replace text in PowerPoint presentation while preserving formatting.
        
        Args:
            ppt_path: Input PowerPoint file path
            output_path: Output PowerPoint file path
            text_data: Original text data with metadata
            translated_texts: List of translated texts
            target_language: Target language code
            
        Returns:
            True if successful, False otherwise
        """
        try:
            prs = Presentation(ppt_path)
            
            # Group translations by slide and shape
            shape_translations = {}
            paragraph_translations = {}
            
            for item, translated_text in zip(text_data, translated_texts):
                slide_idx = item['slide_index']
                shape_idx = item['shape_index']
                
                if item['type'] == 'shape_text':
                    key = (slide_idx, shape_idx)
                    shape_translations[key] = {
                        'text': translated_text,
                        'shape_info': item.get('shape_info', {})
                    }
                elif item['type'] == 'paragraph_text':
                    para_idx = item['paragraph_index']
                    key = (slide_idx, shape_idx, para_idx)
                    paragraph_translations[key] = {
                        'text': translated_text,
                        'paragraph_info': item.get('paragraph_info', {})
                    }
            
            # Apply shape text translations
            for (slide_idx, shape_idx), translation_info in shape_translations.items():
                if (slide_idx < len(prs.slides) and 
                    shape_idx < len(prs.slides[slide_idx].shapes)):
                    
                    shape = prs.slides[slide_idx].shapes[shape_idx]
                    if hasattr(shape, "text"):
                        shape.text = translation_info['text']
                        
                        # Apply formatting adjustments
                        self._apply_shape_format(shape, translation_info['shape_info'], target_language)
                        
                        logger.debug(f"Applied translation to slide {slide_idx + 1}, shape {shape_idx}")
            
            # Apply paragraph translations
            for (slide_idx, shape_idx, para_idx), translation_info in paragraph_translations.items():
                if (slide_idx < len(prs.slides) and 
                    shape_idx < len(prs.slides[slide_idx].shapes)):
                    
                    shape = prs.slides[slide_idx].shapes[shape_idx]
                    if (hasattr(shape, "text_frame") and 
                        para_idx < len(shape.text_frame.paragraphs)):
                        
                        paragraph = shape.text_frame.paragraphs[para_idx]
                        paragraph.text = translation_info['text']
                        
                        # Apply formatting adjustments
                        self._apply_paragraph_format(paragraph, translation_info['paragraph_info'], target_language)
                        
                        logger.debug(f"Applied translation to slide {slide_idx + 1}, shape {shape_idx}, paragraph {para_idx}")
            
            # Save the presentation
            prs.save(output_path)
            
            logger.info(f"Successfully saved translated PowerPoint presentation: {output_path}")
            return True
            
        except Exception as e:
            logger.error(f"Error replacing text in PowerPoint presentation: {e}")
            return False
    
    def _extract_shape_info(self, shape) -> Dict[str, Any]:
        """
        Extract information from a shape.
        
        Args:
            shape: python-pptx shape object
            
        Returns:
            Dictionary containing shape information
        """
        shape_info = {}
        
        try:
            shape_info['shape_type'] = str(shape.shape_type)
            
            if hasattr(shape, 'width'):
                shape_info['width'] = shape.width
            if hasattr(shape, 'height'):
                shape_info['height'] = shape.height
            if hasattr(shape, 'left'):
                shape_info['left'] = shape.left
            if hasattr(shape, 'top'):
                shape_info['top'] = shape.top
            
            # Text frame properties
            if hasattr(shape, 'text_frame'):
                text_frame = shape.text_frame
                if hasattr(text_frame, 'auto_size'):
                    shape_info['auto_size'] = text_frame.auto_size
                if hasattr(text_frame, 'word_wrap'):
                    shape_info['word_wrap'] = text_frame.word_wrap
            
        except Exception as e:
            logger.error(f"Error extracting shape info: {e}")
        
        return shape_info
    
    def _extract_paragraph_info(self, paragraph) -> Dict[str, Any]:
        """
        Extract information from a paragraph.
        
        Args:
            paragraph: python-pptx paragraph object
            
        Returns:
            Dictionary containing paragraph information
        """
        para_info = {}
        
        try:
            if hasattr(paragraph, 'alignment'):
                para_info['alignment'] = paragraph.alignment
            if hasattr(paragraph, 'level'):
                para_info['level'] = paragraph.level
            
            # Font information from runs
            if paragraph.runs:
                first_run = paragraph.runs[0]
                if hasattr(first_run, 'font'):
                    font = first_run.font
                    para_info['font_name'] = font.name
                    para_info['font_size'] = font.size
                    para_info['bold'] = font.bold
                    para_info['italic'] = font.italic
                    para_info['underline'] = font.underline
                    
        except Exception as e:
            logger.error(f"Error extracting paragraph info: {e}")
        
        return para_info
    
    def _apply_shape_format(self, shape, shape_info: Dict[str, Any], target_language: str = "en") -> None:
        """
        Apply formatting adjustments to a shape.
        
        Args:
            shape: python-pptx shape object
            shape_info: Shape information dictionary
            target_language: Target language code
        """
        try:
            # Adjust text frame properties for better text fit
            if hasattr(shape, 'text_frame'):
                text_frame = shape.text_frame
                
                # Enable auto-sizing for better text fit
                if hasattr(text_frame, 'auto_size'):
                    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                
                # Enable word wrap
                if hasattr(text_frame, 'word_wrap'):
                    text_frame.word_wrap = True
                
                # Adjust font for target language
                if target_language == 'th':
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            if hasattr(run, 'font'):
                                run.font.name = 'TH SarabunPSK'
                                if run.font.size:
                                    # Adjust font size
                                    original_size = run.font.size.pt
                                    adjusted_size = max(6, int(original_size * self.font_size_adjustment))
                                    run.font.size = adjusted_size
                
        except Exception as e:
            logger.error(f"Error applying shape format: {e}")
    
    def _apply_paragraph_format(self, paragraph, para_info: Dict[str, Any], target_language: str = "en") -> None:
        """
        Apply formatting adjustments to a paragraph.
        
        Args:
            paragraph: python-pptx paragraph object
            para_info: Paragraph information dictionary
            target_language: Target language code
        """
        try:
            # Restore alignment if available
            if para_info.get('alignment') is not None:
                paragraph.alignment = para_info['alignment']
            
            # Restore level if available
            if para_info.get('level') is not None:
                paragraph.level = para_info['level']
            
            # Apply font adjustments to runs
            for run in paragraph.runs:
                if hasattr(run, 'font'):
                    font = run.font
                    
                    # Font name adjustment for target language
                    if target_language == 'th':
                        font.name = 'TH SarabunPSK'
                    elif para_info.get('font_name'):
                        font.name = para_info['font_name']
                    
                    # Font size adjustment
                    if para_info.get('font_size') and font.size:
                        original_size = font.size.pt
                        adjusted_size = max(6, int(original_size * self.font_size_adjustment))
                        font.size = adjusted_size
                    
                    # Other font properties
                    if para_info.get('bold') is not None:
                        font.bold = para_info['bold']
                    if para_info.get('italic') is not None:
                        font.italic = para_info['italic']
                    if para_info.get('underline') is not None:
                        font.underline = para_info['underline']
                        
        except Exception as e:
            logger.error(f"Error applying paragraph format: {e}")


# Function for simple PowerPoint translation (backward compatibility)
def translate_ppt(input_path: str, output_path: str, target_language: str = "en") -> bool:
    """
    Simple function to translate a PowerPoint presentation.
    
    Args:
        input_path: Path to input PowerPoint file
        output_path: Path to output PowerPoint file
        target_language: Target language code
        
    Returns:
        True if successful, False otherwise
    """
    try:
        processor = PowerPointProcessor()
        return processor.process_file(input_path, output_path, target_language)
    except Exception as e:
        logger.error(f"Error in translate_ppt: {e}")
        return False